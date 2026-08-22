import os
import sys
import argparse
import traceback
import json
import shutil
from datetime import datetime

from pptx import Presentation
# Add root folder to sys.path to resolve local imports
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from job_lock import JobLock
from ipc_messages import IPCMessageSender
from data.data_loader import DataLoader
from data.data_validator import DataValidator
from data.source_snapshot import JobDataSnapshot
from data.revenue_repository import RevenueRepository
from data.inventory_repository import InventoryRepository
from data.models import StoreReportData, StoreMetadata, StaffRoster, StaffItem, OperationalIssue, SECTION_LABELS
from reports.template_preflight import TemplatePreflight
from reports.chart_renderer import ChartRenderer
from reports.image_processor import ImageProcessor
from reports.cluster_aggregator import ClusterAggregator
from reports.pptx_generator import PPTXGenerator
from reports.merge_engine import MergeEngine
from reports.report_qc import ReportQC, QCViolationError

def check_cancellation(cancel_file_path: str, sender: IPCMessageSender) -> bool:
    """Returns True if the GUI created a cancel signal file, False otherwise."""
    if cancel_file_path and os.path.exists(cancel_file_path):
        sender.send_job_cancelled()
        # Clean up the cancel file
        try:
            os.remove(cancel_file_path)
        except Exception:
            pass
        return True
    return False

def send_email_for_store(loader, store_code, store_name, asm_name, report_date, pdf_path, docx_path, pptx_path, xlsx_path, target_email: str = None, cc_emails: str = "", custom_subject: str = "", custom_body: str = "", provider: str = None):
    try:
        from utils.email_dispatcher import EmailDispatcher
        ed = EmailDispatcher(config_dict=loader.config if loader and hasattr(loader, "config") else None)
        email_cfg = ed.get_email_config()
        
        # Resolve target email if empty
        if not target_email and loader and hasattr(loader, "get_asm_email"):
            target_email = loader.get_asm_email(asm_name)
            
        if not target_email:
            print(f"[Email] Không tìm thấy địa chỉ email người nhận cho {store_code}. Bỏ qua.")
            return False

        context = {
            "store_name": store_name or store_code,
            "store_code": store_code,
            "asm_name": asm_name or "QLKD",
            "report_date": report_date or datetime.now().strftime("%d/%m/%Y")
        }

        subject = custom_subject or ed.render_template(email_cfg.get("default_subject_template", "BÁO CÁO KIỂM TRA CỬA HÀNG {store_name} ({store_code}) - NGÀY {report_date}"), context)
        body = custom_body or ed.render_template(email_cfg.get("default_body_template", ""), context)
        cc_final = cc_emails if cc_emails != "" else email_cfg.get("default_cc", "")

        attachments = [p for p in [pptx_path, pdf_path, docx_path, xlsx_path] if p and os.path.exists(p)]

        success, msg = ed.send_report_email(
            to_email=target_email,
            subject=subject,
            body=body,
            cc_emails=cc_final,
            attachments=attachments,
            context=context,
            provider=provider
        )
        print(f"[Email] Result for {store_code}: {msg}")
        return success
    except Exception as e:
        print(f"[Email] Error in send_email_for_store: {e}")
        import traceback
        traceback.print_exc()
        return False
        
    return False

def main():
    parser = argparse.ArgumentParser(description="StoreVisit Subprocess Worker")
    parser.add_argument("--job-id", required=True, help="Unique identifier for the job")
    parser.add_argument("--asm", required=True, help="Selected ASM Name")
    parser.add_argument("--stores", required=True, help="Comma-separated store abbreviations (e.g. BD,VT1,VT2)")
    parser.add_argument("--cancel-file", required=True, help="Path to the cancel signal file")
    parser.add_argument("--form-response-ids", default="", help="Comma-separated list of form response IDs to overlay")
    parser.add_argument("--no-merge", action="store_true", help="Do not merge into a cluster summary. Save separate files per store.")
    parser.add_argument("--schema", default="store_visit", choices=["store_visit", "market_survey", "executive_combo"], help="The schema type to process")
    parser.add_argument("--period-type", default="weekly", choices=["weekly", "monthly", "quarterly"], help="The report period type")
    parser.add_argument("--reference-date", default="auto", help="Date anchor mode: auto, prev, today, or YYYY-MM-DD")
    
    args = parser.parse_args()
    job_id = args.job_id
    asm_name = args.asm
    selected_stores = [s.strip().upper() for s in args.stores.split(",") if s.strip()]
    cancel_file = args.cancel_file
    form_response_ids = [r.strip() for r in args.form_response_ids.split(",") if r.strip()]
    no_merge = args.no_merge
    schema = args.schema
    ref_date = args.reference_date

    # Initialize lock and sender
    lock = JobLock()
    sender = IPCMessageSender(job_id)

    # Acquire lock
    if not lock.acquire():
        sender.send_job_failed("Không thể bắt đầu job: Ứng dụng StoreVisit đang bận chạy một tiến trình khác.")
        sys.exit(1)

    sender.send_job_started({
        "asm": asm_name,
        "stores": selected_stores,
        "cancel_file": cancel_file,
        "form_response_ids": form_response_ids,
        "no_merge": no_merge,
        "schema": schema
    })

    merge_eng = None
    try:
        # Step 1: Load config and verify input files
        sender.send_stage_started("Khởi tạo và cấu hình", total_steps=100)
        sender.send_progress(10, "Đang tải cấu hình...")
        loader = DataLoader()
        
        if schema == "market_survey":
            sender.send_progress(30, "Đang tải dữ liệu khảo sát từ cache...")
            survey_cache_path = loader.config["google"].get("survey_cache_path", "data/survey_cache.json")
            if not os.path.exists(survey_cache_path):
                sender.send_job_failed(f"Không tìm thấy file cache khảo sát tại: {survey_cache_path}. Vui lòng đồng bộ trước.")
                lock.release()
                sys.exit(1)
                
            with open(survey_cache_path, "r", encoding="utf-8") as f:
                survey_cache = json.load(f)
                
            from data.models import MarketSurveyResponse
            responses = []
            mapped_responses = {}
            for rid in form_response_ids:
                if rid in survey_cache:
                    rdata = survey_cache[rid]
                    resp = MarketSurveyResponse(**rdata)
                    responses.append(resp)
                    mapped_responses[rid] = resp
                    
            sender.send_progress(50, f"Đang tổng hợp {len(responses)} phản hồi khảo sát...")
            
            from reports.survey_consolidator import SurveyConsolidator
            consolidator = SurveyConsolidator()
            
            job_temp_dir = os.path.join(loader.get_path("temp_dir"), f"job_{job_id}")
            os.makedirs(job_temp_dir, exist_ok=True)
            temp_excel = os.path.join(job_temp_dir, "survey_report.xlsx")
            
            consolidator.consolidate_to_excel(responses, temp_excel)
            
            sender.send_progress(70, "Đang xuất kết quả khảo sát...")
            output_dir = loader.get_path("output_dir")
            final_job_dir = os.path.join(output_dir, f"BáoCáo_KhảoSát_{job_id}")
            if os.path.exists(final_job_dir):
                shutil.rmtree(final_job_dir)
            os.makedirs(final_job_dir, exist_ok=True)
            
            final_excel = os.path.join(final_job_dir, "survey_report.xlsx")
            shutil.copy2(temp_excel, final_excel)
            
            google_config = loader.config.get("google", {})
            if mapped_responses:
                sender.send_progress(90, "Đang cập nhật trạng thái đã xử lý...")
                sheet_id = google_config.get("spreadsheet_id", "")
                
                reader = None
                if sheet_id:
                    from data.google_sheets_reader import GoogleSheetsReader
                    reader = GoogleSheetsReader(
                        credentials_path=google_config.get("credentials_path", ""),
                        spreadsheet_id=sheet_id
                    )
                
                survey_sheet_name = google_config.get("survey_sheet_name", "MarketSurvey_Responses")
                for rid, resp in mapped_responses.items():
                    survey_cache[rid]["status"] = "done"
                    if reader:
                        try:
                            reader.update_row_status(
                                row_idx=resp.response_id,
                                status="done",
                                sheet_name=survey_sheet_name
                            )
                        except Exception as ex:
                            print(f"Error updating sheet status for response {resp.response_id}: {ex}", file=sys.stderr)
                            survey_cache[rid]["status"] = "error"
                            
                with open(survey_cache_path, "w", encoding="utf-8") as f:
                    json.dump(survey_cache, f, indent=2, ensure_ascii=False)
                    
            sender.send_progress(100, "Đã hoàn thành tổng hợp khảo sát thị trường!")

            
            job_manifest_path = os.path.join(final_job_dir, "manifest.json")
            manifest_data = {
                "job_id": job_id,
                "status": "COMMITTED",
                "timestamp": datetime.now().isoformat(),
                "survey_count": len(responses),
                "files": {
                    "excel": "survey_report.xlsx"
                }
            }
            with open(job_manifest_path, "w", encoding="utf-8") as f:
                json.dump(manifest_data, f, indent=2, ensure_ascii=False)
                
            sender.send_job_completed(final_excel, "", job_manifest_path)
            lock.release()
            sys.exit(0)

        if schema == "executive_combo":
            sender.send_progress(30, f"Đang tổng hợp dữ liệu Báo cáo Executive {args.period_type}...")
            from reports.weekly_monthly_aggregator import WeeklyMonthlyAggregator
            from reports.executive_excel_generator import ExecutiveExcelGenerator
            from reports.executive_pptx_generator import ExecutivePPTXGenerator

            aggregator = WeeklyMonthlyAggregator(loader)
            agg_data = aggregator.aggregate_data(
                period_type=args.period_type,
                asm_filter=asm_name,
                store_filter=selected_stores,
                reference_date=ref_date
            )

            sender.send_progress(60, "Đang tạo Bảng tính Excel Analytics Dashboard (5 Tabs)...")
            output_dir = loader.get_path("output_dir")
            from utils.filename_formatter import format_executive_output_filename
            period_tag = "Tuan" if args.period_type == "weekly" else ("Thang" if args.period_type == "monthly" else "Quy")
            ref_date_str = agg_data.get("period_name", "") or datetime.now().strftime("%d%m%Y")
            
            exec_excel_name = format_executive_output_filename(args.period_type, asm_name, ref_date_str, "xlsx")
            exec_pptx_name = format_executive_output_filename(args.period_type, asm_name, ref_date_str, "pptx")

            final_job_dir = os.path.join(output_dir, f"BaoCao_Executive_{period_tag}_{asm_name}_{job_id[:8]}")
            os.makedirs(final_job_dir, exist_ok=True)

            excel_gen = ExecutiveExcelGenerator()
            excel_path = os.path.join(final_job_dir, exec_excel_name)
            excel_gen.generate(agg_data, excel_path)

            sender.send_progress(85, "Đang tạo Slides PowerPoint Trình Chiếu Executive (5 Slides)...")
            pptx_gen = ExecutivePPTXGenerator()
            pptx_path = os.path.join(final_job_dir, exec_pptx_name)
            pptx_gen.generate(agg_data, pptx_path)

            sender.send_progress(100, f"Đã hoàn thành Báo cáo Combo Executive {args.period_type}!")
            job_manifest_path = os.path.join(final_job_dir, "manifest.json")
            manifest_data = {
                "job_id": job_id,
                "status": "COMMITTED",
                "timestamp": datetime.now().isoformat(),
                "period_type": args.period_type,
                "asm": asm_name,
                "files": {
                    "excel": os.path.basename(excel_path),
                    "pptx": os.path.basename(pptx_path)
                }
            }
            with open(job_manifest_path, "w", encoding="utf-8") as f:
                json.dump(manifest_data, f, indent=2, ensure_ascii=False)

            sender.send_job_completed(excel_path, pptx_path, job_manifest_path)
            lock.release()
            sys.exit(0)

        validator = DataValidator(loader)
        
        # Load form responses if specified
        mapped_responses = {}
        form_cache = None
        if form_response_ids:
            sender.send_progress(15, "Đang tải dữ liệu Google Form từ cache...")
            from data.form_response_cache import FormResponseCache
            cache_path = os.path.join(loader.root_dir, "data", "form_cache.json")
            google_config = loader.config.get("google", {})
            form_cache = FormResponseCache(
                cache_path=cache_path,
                credentials_path=google_config.get("credentials_path", ""),
                spreadsheet_id=google_config.get("spreadsheet_id", ""),
                photo_cache_dir=google_config.get("drive_photos_cache_dir", "")
            )
            for rid in form_response_ids:
                resp = form_cache.get_by_id(rid)
                if resp:
                    mapped_responses[resp.store_code.upper()] = resp
        
        # Snapshot input files before reading
        sender.send_progress(20, "Đang chụp snapshot dữ liệu đầu vào...")
        source_paths = validator.validate_files_exist()
        snapshot = JobDataSnapshot(source_paths)
        snapshot.record_before()

        # Check schemas
        sender.send_progress(30, "Đang xác thực cấu trúc tệp dữ liệu...")
        validator.validate_schemas()
        
        # Preflight templates
        sender.send_progress(40, "Đang kiểm tra chất lượng PowerPoint templates...")
        preflight = TemplatePreflight()
        preflight.verify_templates()
        
        sender.send_stage_completed("Khởi tạo và cấu hình")

        if check_cancellation(cancel_file, sender):
            lock.release()
            sys.exit(0)

        # Step 2: Query and render individual store slides
        sender.send_stage_started("Truy vấn và sinh báo cáo cửa hàng", total_steps=len(selected_stores))
        
        # Repositories
        rev_repo = RevenueRepository(loader, loader.field_mapping)
        inv_repo = InventoryRepository(loader, loader.field_mapping)
        
        # Temp paths for job
        job_temp_dir = os.path.join(loader.get_path("temp_dir"), f"job_{job_id}")
        staging_dir = os.path.join(job_temp_dir, "output_staging")
        os.makedirs(staging_dir, exist_ok=True)
        
        store_reports = []
        store_pptx_paths = []
        
        dim_store = loader.load_dim_store()
        staff_df = loader.load_staff_list()
        weekly_json = loader.load_weekly_json()
        
        store_code_mapping = loader.field_mapping.get("store_code_mapping", {})
        store_staff_workplace = loader.field_mapping.get("store_staff_workplace", {})
        
        for idx, store_key in enumerate(selected_stores):
            if check_cancellation(cancel_file, sender):
                lock.release()
                sys.exit(0)
                
            sender.send_progress(idx, f"Đang xử lý dữ liệu cửa hàng {store_key}...")
            
            store_code = store_code_mapping.get(store_key, store_key)
            
            # Fetch dim store metadata
            row_dim = dim_store[dim_store["StoreCode"] == store_code]
            if row_dim.empty:
                sender.send_warning(f"Không tìm thấy StoreCode '{store_code}' trong DimStore. Bỏ qua.")
                continue
                
            store_name = row_dim.iloc[0]["StoreName"]
            address = row_dim.iloc[0]["Address"]
            region = row_dim.iloc[0]["Region"]
            store_asm = row_dim.iloc[0]["ASM"]
            
            meta = StoreMetadata(
                store_code=store_code,
                store_name=store_name,
                address=address,
                region=region,
                asm_name=str(store_asm) if pd_not_na(store_asm) else asm_name
            )

            # 2.1 Fetch Revenue (Dynamic year/month/cutoff based on report date)
            rev_year = 2026
            rev_month = 7
            cutoff_day = None
            resp = mapped_responses.get(store_key.upper())
            if resp and resp.report_date:
                try:
                    import pandas as pd
                    report_dt = pd.to_datetime(resp.report_date)
                    rev_month = report_dt.month
                    rev_year = report_dt.year
                    cutoff_day = report_dt.day
                except Exception as e:
                    print(f"Error parsing report date {resp.report_date}: {e}")
            
            revenue_data = rev_repo.get_revenue_data(store_key, year=rev_year, month=rev_month, cutoff_day=cutoff_day)
            
            # 2.2 Fetch Inventory & Age
            stock_data = inv_repo.get_stock_inventory(store_key)
            best_sellers = inv_repo.get_best_sellers(store_key)
            slow_sellers = inv_repo.get_slow_sellers(store_key)
            
            # 2.3 Fetch Staff
            workplace_name = store_staff_workplace.get(store_code, store_name)
            df_store_staff = staff_df[staff_df["Nơi làm việc"] == workplace_name]
            
            # Robust fuzzy/normalized workplace matching fallback if direct match empty
            if df_store_staff.empty:
                import re
                def _norm_k(s):
                    if not s: return ""
                    s = str(s).strip()
                    for c1, c2 in [('à|á|ạ|ả|ã|â|ầ|ấ|ậ|ẩ|ẫ|ă|ằ|ắ|ặ|ẳ|ẵ', 'a'),
                                  ('è|é|ẹ|ẻ|ẽ|ê|ề|ế|ệ|ể|ễ', 'e'),
                                  ('ì|í|ị|ỉ|ĩ', 'i'),
                                  ('ò|ó|ọ|ỏ|õ|ô|ồ|ố|ộ|ổ|ỗ|ơ|ờ|ớ|ợ|ở|ỡ', 'o'),
                                  ('ù|ú|ụ|ủ|ũ|ư|ừ|ứ|ự|ử|ữ', 'u'),
                                  ('ỳ|ý|ỵ|ỷ|ỹ', 'y'),
                                  ('đ', 'd'), ('Đ', 'D')]:
                        s = re.sub(f'[{c1}]', c2, s)
                    s = re.sub(r'[^a-zA-Z0-9]', '', s).lower()
                    s = re.sub(r'^0+(\d)', r'\1', s)
                    s = re.sub(r'([a-z])0+(\d)', r'\1\2', s)
                    return s
                
                target_norm = _norm_k(store_name)
                for wp in staff_df["Nơi làm việc"].dropna().unique():
                    wp_norm = _norm_k(wp)
                    if wp_norm == target_norm or wp_norm in target_norm or target_norm in wp_norm:
                        df_store_staff = staff_df[staff_df["Nơi làm việc"] == wp]
                        break

            cht_name = "Chưa bổ nhiệm"
            chp_name = "Chưa bổ nhiệm"
            staff_list = []
            
            for _, r in df_store_staff.iterrows():
                name = r["Tên nhân viên"]
                role = r["Chức danh"] or "NV bán hàng"
                sen = r["Thâm niên năm"]
                sen_val = float(sen) if pd_not_na(sen) else 0.0
                
                if "ch trưởng" in str(role).lower() or "cửa hàng trưởng" in str(role).lower() or str(role).lower() == "cht":
                    cht_name = name
                elif "ch phó" in str(role).lower() or "cửa hàng phó" in str(role).lower() or str(role).lower() == "chp":
                    chp_name = name
                    
                staff_list.append(StaffItem(name=name, role=role, seniority=sen_val))
                
            staff_roster = StaffRoster(cht_name=cht_name, chp_name=chp_name, staff_list=staff_list)
            
            # Fetch form response overlay if available
            form_response = mapped_responses.get(store_code.upper()) or mapped_responses.get(store_key.upper())
            
            # 2.4 Fetch weekly audit note (weekly JSON)
            store_weekly = weekly_json.get(store_key, {})
            
            frontage_rating = "Đạt"
            frontage_issue = "Khu vực trưng bày sạch sẽ, đảm bảo quy chuẩn."
            frontage_action = "Duy trì tiêu chuẩn hàng ngày."
            
            if "sec1" in store_weekly:
                for item in store_weekly["sec1"]:
                    if item.get("label") and "mặt tiền" in str(item["label"]).lower():
                        frontage_rating = item.get("eval") or "Đạt"
                        note_val = item.get("note", "")
                        if note_val:
                            if "->" in str(note_val):
                                frontage_issue, frontage_action = str(note_val).split("->", 1)
                                frontage_issue = frontage_issue.strip()
                                frontage_action = frontage_action.strip()
                            else:
                                frontage_issue = str(note_val)
                                frontage_action = "ASM nhắc nhở khắc phục ngay."
                        break
                        
            # Collect CSVC and issues
            csvc_comment = ""
            issues = []
            issue_idx = 1
            
            # Fetch from checklist JSON if available
            has_checklist_issues = False
            if form_response and form_response.checklist_json:
                try:
                    c_data = json.loads(form_response.checklist_json)
                    sections = c_data.get("sections", {})
                    
                    # Mapping of json section keys to clean labels — nguồn duy nhất ở
                    # data.models.SECTION_LABELS (đồng bộ 30-07: nay có thêm "staff" và
                    # "security_guard", trước đây 2 mục này bị bỏ sót khỏi trích lỗi).
                    sec_labels = SECTION_LABELS

                    for sec_key, label in sec_labels.items():
                        sec_val = sections.get(sec_key, {})
                        for item in sec_val.get("items", []):
                            if item.get("eval") == "Không đạt" and item.get("resolved") != "Có":
                                has_checklist_issues = True
                                item_label = item.get("label", "")
                                note = item.get("note", "")
                                severity = item.get("severity", "Trung bình")
                                assignee = item.get("assignee") or "CHT"
                                deadline = item.get("deadline") or "-"
                                
                                issues.append(OperationalIssue(
                                    index=issue_idx,
                                    label=label,
                                    issue=f"{item_label}: {note}",
                                    date=form_response.report_date or datetime.now().strftime("%d/%m/%Y"),
                                    assignee=assignee,
                                    status="Chưa xử lý",
                                    notes=f"Hạn: {deadline} | Độ ưu tiên: {severity}"
                                ))
                                issue_idx += 1
                                
                    # Set csvc_comment from warehouse/cashier/sub-sections comments if present
                    w_comment_parts = [
                        sections.get("warehouse", {}).get("comment", ""),
                        sections.get("stockroom", {}).get("comment", ""),
                        sections.get("fitting_room", {}).get("comment", ""),
                        sections.get("toilet", {}).get("comment", ""),
                        sections.get("fire_safety", {}).get("comment", "")
                    ]
                    w_comment = " | ".join(filter(None, w_comment_parts))
                    
                    c_comment_parts = [
                        sections.get("cashier", {}).get("comment", ""),
                        sections.get("packaging_security", {}).get("comment", "")
                    ]
                    c_comment = " | ".join(filter(None, c_comment_parts))
                    
                    csvc_comment_parts = []
                    if w_comment:
                        csvc_comment_parts.append(f"Kho: {w_comment}")
                    if c_comment:
                        csvc_comment_parts.append(f"Thu ngân: {c_comment}")
                    if csvc_comment_parts:
                        csvc_comment = " | ".join(csvc_comment_parts)
                except Exception as e:
                    print(f"Error parsing checklist issues: {e}")

            if not has_checklist_issues:
                # Fallback to weekly JSON
                if "sec5" in store_weekly:
                    for item in store_weekly["sec5"]:
                        lbl = item.get("label", "")
                        issue = item.get("issue", "")
                        priority = item.get("priority", "")
                        assignee = item.get("assignee", "")
                        
                        if lbl and "vật chất" in str(lbl).lower() and issue and str(issue).lower() != "ko":
                            csvc_comment = issue
                            
                        if issue and str(issue).lower() != "ko":
                            issues.append(OperationalIssue(
                                index=issue_idx,
                                label=lbl,
                                issue=issue,
                                date=datetime.now().strftime("%d/%m/%Y"),
                                assignee=assignee or "CHT",
                                status="Chưa xử lý",
                                notes=f"Độ ưu tiên: {priority or 'Bình thường'}"
                            ))
                            issue_idx += 1

            # Overlay Form Response values if matched
            if form_response:
                if form_response.asm_name:
                    meta.asm_name = form_response.asm_name
                if form_response.cht_name:
                    staff_roster.cht_name = form_response.cht_name
                if form_response.rating_frontage:
                    frontage_rating = form_response.rating_frontage
                if form_response.comment_frontage:
                    frontage_issue = form_response.comment_frontage
                    frontage_action = form_response.action_plan or "ASM nhắc nhở khắc phục ngay."
                if form_response.comment_csvc and not csvc_comment:
                    csvc_comment = form_response.comment_csvc
                if form_response.pending_issues and not has_checklist_issues:
                    # Only append old summary pending issue if we didn't extract detailed ones
                    issues.append(OperationalIssue(
                        index=issue_idx,
                        label="Đề xuất từ Form",
                        issue=form_response.pending_issues,
                        date=form_response.report_date or datetime.now().strftime("%d/%m/%Y"),
                        assignee="CHT",
                        status="Chưa xử lý",
                        notes=f"Hạn: {form_response.action_deadline or ''} | Kế hoạch: {form_response.action_plan or ''}"
                    ))
                    issue_idx += 1

            store_report = StoreReportData(
                metadata=meta,
                revenue=revenue_data,
                stock=stock_data,
                best_sellers=best_sellers,
                slow_sellers=slow_sellers,
                staff=staff_roster,
                issues=issues,
                csvc_comment=csvc_comment,
                frontage_rating=frontage_rating,
                frontage_issue=frontage_issue,
                frontage_action=frontage_action,
                form_response=form_response
            )
            store_reports.append(store_report)

            # 2.5 Generate helper assets (chart image, frontage photos)
            chart_path = os.path.join(job_temp_dir, f"chart_{store_code}.png")
            chart_rend = ChartRenderer()
            chart_rend.draw_revenue_chart(
                actual=revenue_data.revenue_actual,
                target=revenue_data.revenue_target,
                prev=revenue_data.revenue_prev,
                yoy=revenue_data.revenue_yoy,
                output_path=chart_path
            )
            
            # Process and fit all photos in form_response in-place (orientation/compression only)
            # If local_path is missing or file not found, re-download from Drive
            if form_response:
                google_config = loader.config.get("google", {})
                _drive_dl = None
                _photo_cache_dir = google_config.get("drive_photos_cache_dir", "")
                if not _photo_cache_dir:
                    _photo_cache_dir = os.path.join(loader.root_dir, "temp", "photo_cache")
                os.makedirs(_photo_cache_dir, exist_ok=True)

                for photo in form_response.photos:
                    # Step 1: Re-download if local_path is absent or stale
                    if not photo.local_path or not os.path.exists(photo.local_path):
                        if photo.drive_url:
                            try:
                                if _drive_dl is None:
                                    from data.drive_image_downloader import DriveImageDownloader
                                    _drive_dl = DriveImageDownloader(google_config.get("credentials_path", ""))
                                local_p = _drive_dl.download_image(photo.drive_url, _photo_cache_dir)
                                if local_p:
                                    photo.local_path = local_p
                                    print(f"[Photo] Re-downloaded {photo.section}[{photo.index}]: {local_p}")
                                else:
                                    print(f"[Photo] Download failed for {photo.section}[{photo.index}]: {photo.drive_url}")
                            except Exception as _dl_err:
                                print(f"[Photo] Error re-downloading {photo.section}[{photo.index}]: {_dl_err}")

                    # Step 2: Compress + orient the photo into the job temp dir
                    if photo.local_path and os.path.exists(photo.local_path):
                        dest_path = os.path.join(job_temp_dir, f"photo_{photo.section}_{photo.index}_{store_code}.jpg")
                        ImageProcessor.process_and_compress_image(photo.local_path, dest_path)
                        photo.local_path = dest_path

            # Map frontage photos from form_response
            front_photos = []
            if form_response:
                front_photos = [p for p in form_response.photos if p.section == "frontage" and p.local_path]
                
            # Check if revenue data was missing (i.e. no records at all in revenue or target)
            df_rev = loader.load_revenue()
            df_tgt = loader.load_target()
            store_code_upper = store_code.upper()
            has_rev_records = (df_rev["StoreCode"].str.upper() == store_code_upper).any()
            has_tgt_records = (df_tgt["StoreCode"].str.upper() == store_code_upper).any()
            
            # Check if stock data was missing
            df_stock = loader.load_stock()
            has_stock_records = (df_stock["StoreCode"].str.upper() == store_code_upper).any()

            # 2.5b Generate stock doughnut chart
            stock_chart_path = os.path.join(job_temp_dir, f"stock_chart_{store_code}.png")
            if has_stock_records and stock_data:
                chart_rend.draw_stock_doughnut_chart(
                    age_groups=stock_data.age_groups,
                    total_skus=stock_data.skus_count,
                    output_path=stock_chart_path
                )
            else:
                stock_chart_path = ""

            img_paths = {
                "revenue_chart": chart_path,
                "stock_chart": stock_chart_path,
                "front_left": front_photos[0].local_path if len(front_photos) > 0 else "",
                "front_center": front_photos[1].local_path if len(front_photos) > 1 else "",
                "front_right": front_photos[2].local_path if len(front_photos) > 2 else "",
                "missing_revenue": not (has_rev_records or has_tgt_records),
                "missing_stock": not has_stock_records
            }

            # 2.6 Generate PPTX
            store_pptx = os.path.join(job_temp_dir, f"report_{store_code}.pptx")
            master_store_template = os.path.join(loader.get_path("templates_dir"), "Store_Report_Master.pptx")
            
            gen = PPTXGenerator(master_store_template)
            gen.generate_store_report(store_report, store_pptx, img_paths)
            
            # Generate DOCX report
            try:
                from reports.docx_generator import DocxGenerator
                docx_gen = DocxGenerator(job_temp_dir)
                store_docx = docx_gen.generate(store_report)
                # Ensure correct name format
                std_docx = os.path.join(job_temp_dir, f"report_{store_code}.docx")
                if os.path.exists(std_docx):
                    try: os.remove(std_docx)
                    except: pass
                os.rename(store_docx, std_docx)
                print(f"[Worker] Generated DOCX report successfully at: {std_docx}")
            except Exception as docx_err:
                print(f"[Worker] Error generating DOCX report for {store_code}: {docx_err}")
                traceback.print_exc()

            # Generate XLSX report
            try:
                from reports.excel_generator import ExcelGenerator
                excel_gen = ExcelGenerator(job_temp_dir)
                store_xlsx = excel_gen.generate(store_report)
                # Ensure correct name format
                std_xlsx = os.path.join(job_temp_dir, f"report_{store_code}.xlsx")
                if os.path.exists(std_xlsx):
                    try: os.remove(std_xlsx)
                    except: pass
                os.rename(store_xlsx, std_xlsx)
                print(f"[Worker] Generated XLSX report successfully at: {std_xlsx}")
            except Exception as excel_err:
                print(f"[Worker] Error generating XLSX report for {store_code}: {excel_err}")
                traceback.print_exc()
            
            store_pptx_paths.append((store_key, store_pptx))
            
        sender.send_stage_completed("Truy vấn và sinh báo cáo cửa hàng")

        if check_cancellation(cancel_file, sender):
            lock.release()
            sys.exit(0)

        if no_merge:
            # Step 3: Export separate presentations to PDF via PowerPoint COM
            sender.send_stage_started("Xuất bản báo cáo riêng lẻ và PDF", total_steps=100)
            sender.send_progress(10, "Khởi tạo tiến trình PowerPoint COM...")
            merge_eng = MergeEngine()
            merge_eng.start_powerpoint()
            
            for i, (store_key, store_pptx) in enumerate(store_pptx_paths):
                pct = 20 + int(70 * (i / len(store_pptx_paths)))
                store_code = store_code_mapping.get(store_key, store_key)
                sender.send_progress(pct, f"Đang xuất PDF cho cửa hàng {store_code}...")
                
                abs_pptx = os.path.abspath(store_pptx)
                pdf_path = os.path.join(job_temp_dir, f"report_{store_code}.pdf")
                abs_pdf = os.path.abspath(pdf_path)
                
                try:
                    deck = merge_eng.powerpoint.Presentations.Open(abs_pptx)
                    deck.SaveAs(abs_pdf, 32) # 32 is ppSaveAsPDF
                    deck.Close()
                except Exception as ex:
                    print(f"Error exporting PDF for store {store_code}: {ex}", file=sys.stderr)
                    
            sender.send_stage_completed("Xuất bản báo cáo riêng lẻ và PDF")
            
            if check_cancellation(cancel_file, sender):
                lock.release()
                if merge_eng:
                    merge_eng.close_powerpoint()
                sys.exit(0)
                
            # Step 4: Quality Control & Atomic promotion for separate reports
            sender.send_stage_started("Kiểm định chất lượng (QC)", total_steps=100)
            sender.send_progress(20, "Đang chạy bộ lọc QC kiểm tra chất lượng tệp...")
            
            qc = ReportQC()
            # Verify each store's report
            # Store master has 19 slides
            for store_key, store_pptx in store_pptx_paths:
                store_code = store_code_mapping.get(store_key, store_key)
                pdf_path = os.path.join(job_temp_dir, f"report_{store_code}.pdf")
                dynamic_slide_count = len(Presentation(store_pptx).slides)
                qc.verify_report(store_pptx, pdf_path, expected_slide_count=dynamic_slide_count, com_engine=merge_eng)
                
            sender.send_progress(50, "Xác minh độ toàn vẹn của tệp nguồn...")
            snapshot.verify_after_and_save(os.path.join(staging_dir, "source_snapshot.json"))
            
            # Close COM before renaming to avoid locks
            sender.send_progress(70, "Giải phóng tiến trình PowerPoint COM...")
            merge_eng.close_powerpoint()
            merge_eng = None
            
            sender.send_progress(90, "Đang chuyển giao tệp kết quả chính thức...")
            output_dir = loader.get_path("output_dir")
            
            # Create final output folder
            final_job_dir = os.path.join(output_dir, f"BaoCao_Rieng_{job_id}")
            if os.path.exists(final_job_dir):
                shutil.rmtree(final_job_dir)
            os.makedirs(final_job_dir, exist_ok=True)
            
            # Copy each store's file into a subfolder
            from utils.filename_formatter import format_store_output_filename
            final_job_dirs = []
            for store_rep in store_reports:
                s_code = store_rep.metadata.store_code
                s_name = store_rep.metadata.store_name
                s_asm = store_rep.metadata.asm_name
                s_date = store_rep.form_response.report_date if store_rep.form_response else ""
                
                store_subfolder = os.path.join(final_job_dir, f"BaoCao_{s_code}")
                os.makedirs(store_subfolder, exist_ok=True)
                
                fn_pptx = format_store_output_filename(s_name, s_asm, s_date, "pptx")
                fn_pdf = format_store_output_filename(s_name, s_asm, s_date, "pdf")
                fn_docx = format_store_output_filename(s_name, s_asm, s_date, "docx")
                fn_xlsx = format_store_output_filename(s_name, s_asm, s_date, "xlsx")
                
                pptx_src = os.path.join(job_temp_dir, f"report_{s_code}.pptx")
                pdf_src = os.path.join(job_temp_dir, f"report_{s_code}.pdf")
                docx_src = os.path.join(job_temp_dir, f"report_{s_code}.docx")
                xlsx_src = os.path.join(job_temp_dir, f"report_{s_code}.xlsx")
                
                pptx_p = os.path.join(store_subfolder, fn_pptx)
                pdf_p = os.path.join(store_subfolder, fn_pdf)
                docx_p = os.path.join(store_subfolder, fn_docx)
                xlsx_p = os.path.join(store_subfolder, fn_xlsx)
                
                if os.path.exists(pptx_src): shutil.copy2(pptx_src, pptx_p)
                if os.path.exists(pdf_src): shutil.copy2(pdf_src, pdf_p)
                if os.path.exists(docx_src): shutil.copy2(docx_src, docx_p)
                if os.path.exists(xlsx_src): shutil.copy2(xlsx_src, xlsx_p)
                
                # Trigger email sending (Disabled by default to prevent unwanted background emails)
                ENABLE_AUTO_EMAIL = False
                if ENABLE_AUTO_EMAIL:
                    send_email_for_store(
                        loader=loader,
                        store_code=store_code,
                        store_name=store_name,
                        asm_name=asm_name_val,
                        report_date=report_date_val,
                        pdf_path=pdf_p,
                        docx_path=docx_p,
                        pptx_path=pptx_p,
                        xlsx_path=xlsx_p
                    )
                else:
                    print(f"[Email] Tự động gửi mail đang TẮT (ENABLE_AUTO_EMAIL=False). Bỏ qua gửi mail cho {store_code}.")

                # Write individual manifest
                manifest_data = {
                    "job_id": job_id,
                    "store_code": store_code,
                    "status": "COMMITTED",
                    "timestamp": datetime.now().isoformat(),
                    "files": {
                        "pptx": "report.pptx",
                        "pdf": "report.pdf",
                        "docx": "report.docx",
                        "xlsx": "report.xlsx"
                    }
                }
                with open(os.path.join(store_subfolder, "manifest.json"), "w", encoding="utf-8") as f:
                    json.dump(manifest_data, f, indent=2, ensure_ascii=False)
                final_job_dirs.append(store_subfolder)
            
            # Write overall job manifest
            manifest_data = {
                "job_id": job_id,
                "status": "COMMITTED",
                "timestamp": datetime.now().isoformat(),
                "store_count": len(selected_stores),
                "merge_mode": "separate",
                "files": {
                    "output_folder": final_job_dir
                }
            }
            with open(os.path.join(final_job_dir, "manifest.json"), "w", encoding="utf-8") as f:
                json.dump(manifest_data, f, indent=2, ensure_ascii=False)
                
            # Mark used responses as done in cache and Sheets
            if form_cache and mapped_responses:
                sender.send_progress(95, "Đang cập nhật trạng thái đã xử lý lên Google Sheets...")
                for rid, resp in mapped_responses.items():
                    try:
                        form_cache.mark_done(resp.response_id)
                    except Exception as ex:
                        print(f"Error marking response {resp.response_id} as done: {ex}", file=sys.stderr)
                        
            sender.send_stage_completed("Kiểm định chất lượng (QC)")
            
            if not store_pptx_paths:
                raise ValueError("Không có cửa hàng nào được tạo báo cáo thành công do lỗi dữ liệu đầu vào hoặc không map được Store Code.")
                
            first_rep = store_reports[0]
            first_code = first_rep.metadata.store_code
            first_name = first_rep.metadata.store_name
            first_asm = first_rep.metadata.asm_name
            first_date = first_rep.form_response.report_date if first_rep.form_response else ""
            
            f_pptx_name = format_store_output_filename(first_name, first_asm, first_date, "pptx")
            f_pdf_name = format_store_output_filename(first_name, first_asm, first_date, "pdf")
            
            final_pptx = os.path.join(final_job_dir, f"BaoCao_{first_code}", f_pptx_name)
            final_pdf = os.path.join(final_job_dir, f"BaoCao_{first_code}", f_pdf_name)
            final_manifest = os.path.join(final_job_dir, "manifest.json")
            sender.send_job_completed(final_pptx, final_pdf, final_manifest)

        else:
            # Step 3: Aggregate and Merge slides
            sender.send_stage_started("Gộp báo cáo cụm và xuất PDF", total_steps=100)
            sender.send_progress(10, "Đang tính toán cộng dồn KPIs cụm...")
            agg = ClusterAggregator()
            cluster_data = agg.aggregate_stores(store_reports, cluster_name=f"Cụm {asm_name}")
            
            sender.send_progress(30, "Khởi tạo tiến trình PowerPoint COM...")
            merge_eng = MergeEngine()
            merge_eng.start_powerpoint()
            
            temp_merged_pptx = os.path.join(staging_dir, "report.pptx")
            temp_merged_pdf = os.path.join(staging_dir, "report.pdf")
            
            sender.send_progress(50, "Đang ghép các slide báo cáo...")
            merge_eng.merge_and_export(
                cluster_data=cluster_data,
                store_pptx_paths=store_pptx_paths,
                final_pptx_path=temp_merged_pptx,
                final_pdf_path=temp_merged_pdf
            )
            
            sender.send_stage_completed("Gộp báo cáo cụm và xuất PDF")

            if check_cancellation(cancel_file, sender):
                lock.release()
                if merge_eng:
                    merge_eng.close_powerpoint()
                sys.exit(0)

            # Step 4: Quality Control & Atomic folder promotion
            sender.send_stage_started("Kiểm định chất lượng (QC)", total_steps=100)
            sender.send_progress(20, "Đang chạy bộ lọc QC kiểm tra chất lượng tệp...")
            
            expected_slide_count = 2 + sum(
                len([s for s in Presentation(path).slides if merge_eng._get_slide_id(s) not in merge_eng.exclude_slide_ids])
                for _, path in store_pptx_paths
            )
            
            qc = ReportQC()
            qc.verify_report(temp_merged_pptx, temp_merged_pdf, expected_slide_count, com_engine=merge_eng)
            
            # Verification that source files haven't changed
            sender.send_progress(50, "Xác minh độ toàn vẹn của tệp nguồn...")
            job_manifest_path = os.path.join(staging_dir, "manifest.json")
            snapshot.verify_after_and_save(os.path.join(staging_dir, "source_snapshot.json"))
            
            # Close COM before renaming to avoid PowerPoint locking the staging files!
            sender.send_progress(70, "Giải phóng tiến trình PowerPoint COM...")
            merge_eng.close_powerpoint()
            merge_eng = None
            
            from utils.filename_formatter import format_cluster_output_filename
            cluster_name_val = f"Cum_{asm_name}"
            cluster_pptx_name = format_cluster_output_filename(cluster_name_val, asm_name, ref_date, "pptx")
            cluster_pdf_name = format_cluster_output_filename(cluster_name_val, asm_name, ref_date, "pdf")

            # Create manifest
            manifest_data = {
                "job_id": job_id,
                "status": "COMMITTED",
                "timestamp": datetime.now().isoformat(),
                "store_count": len(selected_stores),
                "expected_slides": expected_slide_count,
                "files": {
                    "pptx": cluster_pptx_name,
                    "pdf": cluster_pdf_name
                }
            }
            with open(job_manifest_path, "w", encoding="utf-8") as f:
                json.dump(manifest_data, f, indent=2, ensure_ascii=False)
                
            # Promote folder atomicly
            sender.send_progress(90, "Đang chuyển giao tệp kết quả chính thức...")
            output_dir = loader.get_path("output_dir")
            final_job_dir = os.path.join(output_dir, f"BaoCao_Cum_{job_id}")
            
            if os.path.exists(final_job_dir):
                shutil.rmtree(final_job_dir)
                
            # Mark used responses as done in cache and Sheets
            if form_cache and mapped_responses:
                sender.send_progress(95, "Đang cập nhật trạng thái đã xử lý lên Google Sheets...")
                for rid, resp in mapped_responses.items():
                    try:
                        form_cache.mark_done(resp.response_id)
                    except Exception as ex:
                        print(f"Error marking response {resp.response_id} as done: {ex}", file=sys.stderr)

            shutil.move(staging_dir, final_job_dir)
            
            # Rename staging files to standard cluster filenames inside final_job_dir
            old_pptx = os.path.join(final_job_dir, "report.pptx")
            old_pdf = os.path.join(final_job_dir, "report.pdf")
            final_pptx = os.path.join(final_job_dir, cluster_pptx_name)
            final_pdf = os.path.join(final_job_dir, cluster_pdf_name)
            
            if os.path.exists(old_pptx):
                if os.path.exists(final_pptx): os.remove(final_pptx)
                os.rename(old_pptx, final_pptx)
            if os.path.exists(old_pdf):
                if os.path.exists(final_pdf): os.remove(final_pdf)
                os.rename(old_pdf, final_pdf)

            sender.send_stage_completed("Kiểm định chất lượng (QC)")
            
            final_manifest = os.path.join(final_job_dir, "manifest.json")
            sender.send_job_completed(final_pptx, final_pdf, final_manifest)
        
    except Exception as e:
        tb = traceback.format_exc()
        print(f"Error occurred in job worker: {e}\n{tb}", file=sys.stderr)
        sender.send_job_failed(str(e), tb)
        
    finally:
        # Cleanup
        if merge_eng:
            try:
                merge_eng.close_powerpoint()
            except Exception:
                pass
        lock.release()

def pd_not_na(val) -> bool:
    try:
        import pandas as pd
        return pd.notna(val)
    except Exception:
        return val is not None

if __name__ == "__main__":
    main()
