import os
import time
from datetime import datetime
import win32com.client
import win32process
import win32gui
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from data.models import ClusterReportData

class MergeEngine:
    def __init__(self, config_path: str = "config/app_config.yaml"):
        self.root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.config_path = os.path.join(self.root_dir, config_path)
        with open(self.config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)
            
        self.exclude_slide_ids = self.config["cluster_merge"]["exclude_store_slide_ids"]
        self.powerpoint = None
        self.ppt_hwnd = None
        self.ppt_pid = None

    def start_powerpoint(self):
        """Start a private PowerPoint COM session and record its PID/HWND for safe termination."""
        print("Starting private PowerPoint COM session...")
        # Get list of running PIDs before DispatchEx
        pids_before = self._get_powerpoint_pids()
        
        # Spawn application
        self.powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        self.powerpoint.Visible = 1
        
        # Get window handle safely
        try:
            self.ppt_hwnd = self.powerpoint.HWND()
        except Exception:
            try:
                self.ppt_hwnd = win32gui.FindWindow("PP12FrameClass", None)
            except Exception:
                self.ppt_hwnd = None
        
        # Get actual PID from HWND
        if self.ppt_hwnd:
            try:
                _, pid = win32process.GetWindowThreadProcessId(self.ppt_hwnd)
                self.ppt_pid = pid
                print(f"PowerPoint started with HWND: {self.ppt_hwnd}, PID: {self.ppt_pid}")
            except Exception:
                pass
        
        # Fallback to PID diff if PID could not be determined from HWND
        if not self.ppt_pid:
            pids_after = self._get_powerpoint_pids()
            diff = pids_after - pids_before
            if diff:
                self.ppt_pid = list(diff)[0]
                print(f"PowerPoint started (PID diff): {self.ppt_pid}")

    def merge_and_export(self, cluster_data: ClusterReportData, store_pptx_paths: list, final_pptx_path: str, final_pdf_path: str):
        """
        Merge individual store presentations into the cluster summary deck.
        Saves the merged PPTX and exports to PDF.
        """
        # Ensure PowerPoint is started
        if not self.powerpoint:
            self.start_powerpoint()

        # 1. Create a copy of the Cluster Summary Master template
        templates_dir = self.config["paths"]["templates_dir"]
        if not os.path.isabs(templates_dir):
            templates_dir = os.path.join(self.root_dir, templates_dir)
            
        cluster_master_path = os.path.join(templates_dir, "Cluster_Summary_Master.pptx")
        
        # Generate the cluster summary slide text first via python-pptx
        prs = Presentation(cluster_master_path)
        self._populate_cluster_slides(prs, cluster_data)
        prs.save(final_pptx_path)

        # 2. Open final presentation via COM
        abs_final_pptx = os.path.abspath(final_pptx_path)
        deck = self.powerpoint.Presentations.Open(abs_final_pptx)
        
        # 3. Insert slides from each store presentation, excluding Cover, ToC, and Thank You slides
        current_slide_count = deck.Slides.Count
        
        for store_key, store_path in store_pptx_paths:
            abs_store_path = os.path.abspath(store_path)
            
            # Find slide indices in individual pptx that should be copied
            # We use python-pptx to read the slide metadata in memory quickly
            store_prs = Presentation(store_path)
            slides_to_copy = [] # 1-indexed for PowerPoint COM
            
            for i, slide in enumerate(store_prs.slides):
                slide_id = self._get_slide_id(slide)
                if slide_id not in self.exclude_slide_ids:
                    slides_to_copy.append(i + 1)
                    
            # Group slides into contiguous ranges to speed up COM insertion
            ranges = self._get_contiguous_ranges(slides_to_copy)
            
            # Insert ranges
            for start, end in ranges:
                # InsertFromFile(FileName, Index, SlideStart, SlideEnd)
                # Index is where to insert (current_slide_count inserts at the end)
                deck.Slides.InsertFromFile(abs_store_path, current_slide_count, start, end)
                current_slide_count = deck.Slides.Count

        # Save merged presentation
        deck.Save()
        print(f"Merged presentation saved to {final_pptx_path}")

        # 4. Export to PDF
        abs_final_pdf = os.path.abspath(final_pdf_path)
        print(f"Exporting presentation to PDF at {final_pdf_path}...")
        deck.SaveAs(abs_final_pdf, 32) # 32 is ppSaveAsPDF
        deck.Close()
        print("PDF Export completed successfully.")

    def close_powerpoint(self):
        """Cleanly quit the PowerPoint application object."""
        if self.powerpoint:
            try:
                self.powerpoint.Quit()
            except Exception as e:
                print(f"Warning: Exception while quitting PowerPoint COM: {e}")
            self.powerpoint = None
            
        # Verify the process is dead, if not terminate specifically
        if self.ppt_pid:
            self._ensure_pid_terminated(self.ppt_pid)

    def _populate_cluster_slides(self, prs, data: ClusterReportData):
        """Populate the cover and table slide of the cluster report."""
        # ── Slide 0: CLUSTER_COVER ──
        # Named shapes TXT_CLUSTER_NAME, TXT_REPORT_DATE, TXT_ASM_NAME
        # were renamed by rebuild_templates.py from the original template's visible shapes.
        cover_slide = prs.slides[0]
        asm_name = data.individual_reports[0].metadata.asm_name if data.individual_reports else "QLKD"

        self._fill_text(cover_slide, "TXT_CLUSTER_NAME",
                        f"{data.cluster_name}", font_size=24, bold=True)
        self._fill_text(cover_slide, "TXT_REPORT_DATE",
                        data.report_date, font_size=14)
        self._fill_text(cover_slide, "TXT_ASM_NAME",
                        asm_name, font_size=14)
        # CHT field is not applicable for cluster — clear it
        self._fill_text(cover_slide, "TXT_CHT_NAME", "", font_size=14)

        # Add KPI summary boxes dynamically (positioned in the right panel of the Cover)
        # Original template right panel area: left ≈ 8.4", top ≈ 4.3"-9.0"
        kpi_box_specs = [
            ("KPI_CLUSTER_REVENUE_ACTUAL",
             f"Doanh thu: {data.revenue_actual / 1_000_000_000:.2f} Tỷ đồng",
             Inches(8.4), Inches(8.0), Inches(10.5), Inches(0.7),
             16, True, RGBColor(10, 35, 66)),
            ("KPI_CLUSTER_REVENUE_ATTAINMENT",
             f"Đạt kế hoạch: {data.attainment_pct:.1f}%",
             Inches(8.4), Inches(8.85), Inches(10.5), Inches(0.6),
             14, False, RGBColor(0, 120, 60) if data.attainment_pct >= 100 else RGBColor(180, 30, 30)),
        ]
        existing_shape_names = {s.name for s in cover_slide.shapes}
        for (name, text, left, top, width, height, fsize, bold, color) in kpi_box_specs:
            if name not in existing_shape_names:
                tb = cover_slide.shapes.add_textbox(left, top, width, height)
                tb.name = name
            else:
                tb = next(s for s in cover_slide.shapes if s.name == name)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.text = text
            if p.runs:
                run = p.runs[0]
                run.font.name = "Inter"
                run.font.size = Pt(fsize)
                run.font.bold = bold
                run.font.color.rgb = color


        # 2. Slide 1: CLUSTER_SUMMARY
        summary_slide = prs.slides[1]
        
        # Delete the original dummy shape grid representing the table to prevent overlap
        shapes_to_delete = []
        for s in summary_slide.shapes:
            # Delete shape elements in the table area (Top > 2.2 inches)
            # Do not delete metadata slide ID shape or footnote
            if s.top > Inches(2.2) and s.name not in ["META_SLIDE_ID", "META_TEMPLATE_VERSION"] and not (s.has_text_frame and "🟢" in s.text_frame.text):
                shapes_to_delete.append(s)
                
        for s in shapes_to_delete:
            summary_slide.shapes._spTree.remove(s._element)

        # Draw a beautiful clean table for performance
        rows_data = []
        for item in data.stores_performance:
            rows_data.append([
                item.store_name,
                f"{item.revenue_actual / 1_000_000_000:.2f} Tỷ",
                f"{item.revenue_target / 1_000_000_000:.2f} Tỷ",
                f"{item.attainment_pct:.1f}%"
            ])
            
        headers = ["Cửa hàng", "Doanh thu (Lũy kế)", "Kế hoạch tháng", "Đạt kế hoạch %"]
        col_widths = [Inches(5.0), Inches(4.5), Inches(4.5), Inches(4.5)]
        aligns = ["left", "right", "right", "right"]
        
        # Insert table
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(19.0)
        height = Inches(6.0)
        
        table_shape = summary_slide.shapes.add_table(len(rows_data) + 1, len(headers), left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        for col_idx, w in enumerate(col_widths):
            table.columns[col_idx].width = w
            
        # Write headers
        header_bg = RGBColor(10, 35, 66)
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            self._format_cell(cell, h, font_size=12, bold=True, color=RGBColor(255, 255, 255), align="center", bg_color=header_bg)
            
        # Write rows
        for r_idx, row_vals in enumerate(rows_data):
            row_bg = RGBColor(255, 255, 255) if r_idx % 2 == 0 else RGBColor(245, 248, 252)
            for c_idx, val in enumerate(row_vals):
                cell = table.cell(r_idx + 1, c_idx)
                align = aligns[c_idx]
                bold = True if c_idx == 0 else False
                self._format_cell(cell, val, font_size=11, bold=bold, color=RGBColor(0, 0, 0), align=align, bg_color=row_bg)

    def _format_cell(self, cell, text: str, font_size: int, bold: bool, color: RGBColor, align: str, bg_color: RGBColor):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.name = "Inter"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
            
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    def _fill_text(self, slide, shape_name: str, new_text: str, font_size: int = 14, bold: bool = False, align: str = "left"):
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                font_name = "Inter"
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    if run.font.name:
                        font_name = run.font.name
                
                shape.text_frame.text = ""
                p = shape.text_frame.paragraphs[0]
                p.text = str(new_text)
                
                if align == "center":
                    p.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
                if p.runs:
                    run = p.runs[0]
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                return

    def _get_slide_id(self, slide) -> str:
        for shape in slide.shapes:
            if shape.name == "META_SLIDE_ID" and shape.has_text_frame:
                return shape.text_frame.text.strip()
        return ""

    def _get_contiguous_ranges(self, indices: list) -> list:
        """Group a list of integers into contiguous ranges. e.g. [2, 3, 4, 6] -> [(2, 4), (6, 6)]"""
        if not indices:
            return []
        indices.sort()
        ranges = []
        start = indices[0]
        end = indices[0]
        for x in indices[1:]:
            if x == end + 1:
                end = x
            else:
                ranges.append((start, end))
                start = x
                end = x
        ranges.append((start, end))
        return ranges

    def _get_powerpoint_pids(self) -> set:
        """Get set of PIDs of currently running powerpoint processes."""
        import subprocess
        pids = set()
        try:
            output = subprocess.check_output(
                'tasklist /FI "IMAGENAME eq POWERPNT.EXE" /FO CSV /NH',
                shell=True, stderr=subprocess.DEVNULL
            ).decode("utf-8", errors="ignore")
            for line in output.splitlines():
                if line.strip():
                    parts = line.split(",")
                    if len(parts) > 1:
                        pid_str = parts[1].strip('"')
                        if pid_str.isdigit():
                            pids.add(int(pid_str))
        except Exception:
            pass
        return pids

    def _ensure_pid_terminated(self, pid: int):
        """Force kill the specific PID if it didn't shut down cleanly."""
        import subprocess
        # Verify it is still running
        pids = self._get_powerpoint_pids()
        if pid in pids:
            print(f"Warning: PowerPoint process (PID: {pid}) did not quit cleanly. Terminating process...")
            try:
                subprocess.run(f"taskkill /PID {pid} /F", shell=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                print(f"Process PID: {pid} killed successfully.")
            except Exception as e:
                print(f"Error killing process PID: {pid}: {e}")
