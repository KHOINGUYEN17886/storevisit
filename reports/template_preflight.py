import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PreflightError(Exception):
    pass

class TemplatePreflight:
    def __init__(self, config_path: str = "config/app_config.yaml", manifest_path: str = "config/template_manifest.yaml"):
        # Resolve paths
        self.root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.config_path = os.path.join(self.root_dir, config_path)
        self.manifest_path = os.path.join(self.root_dir, manifest_path)
        
        with open(self.config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)
            
        with open(self.manifest_path, "r", encoding="utf-8") as f:
            self.manifest = yaml.safe_load(f)
            
        self.templates_dir = self.config["paths"]["templates_dir"]
        if not os.path.isabs(self.templates_dir):
            self.templates_dir = os.path.join(self.root_dir, self.templates_dir)

    def initialize_metadata_and_shapes(self):
        """
        One-time setup to inject META_SLIDE_ID and rename shapes in raw templates.
        This runs during installation or setup.
        """
        print("Initializing templates metadata and shape names...")
        
        # 1. Setup Store Master
        store_path = os.path.join(self.templates_dir, self.manifest["templates"]["store_master"])
        if not os.path.exists(store_path):
            raise FileNotFoundError(f"Store master template not found at {store_path}")
            
        prs = Presentation(store_path)
        
        # Mapping slide indices to Slide IDs (covering all processed slides)
        store_slide_map = {
            0: "STORE_COVER",
            1: "STORE_GENERAL_INFO",
            2: "STORE_FRONTAGE_PHOTOS",
            3: "STORE_INNER_PHOTOS",
            4: "STORE_VM_ERROR_1",
            5: "STORE_VM_ERROR_2",
            6: "STORE_VM_ERROR_3",
            7: "STORE_VM_ERROR_4",
            8: "STORE_STOCKROOM_CASHIER",
            9: "STORE_COMPETITOR",
            10: "STORE_REVENUE",
            11: "STORE_STOCK_INVENTORY",
            12: "STORE_BEST_SELLERS",
            13: "STORE_SLOW_SELLERS",
            14: "STORE_STAFF_LIST",
            15: "STORE_SURVEY_SUPPORT",
            16: "STORE_PENDING_ISSUES",
            17: "STORE_DEV_PROPOSALS",
            18: "STORE_THANK_YOU"
        }
        
        # Rename shapes and inject META_SLIDE_ID
        for idx, slide_id in store_slide_map.items():
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]
            
            # Inject metadata shapes (idempotent)
            self._inject_metadata_shape(slide, slide_id)
            
            # Rename shapes on specific slides
            if slide_id == "STORE_COVER":
                self._rename_shape(slide, "Text 5", "TXT_STORE_NAME")
                self._rename_shape(slide, "Text 7", "TXT_REPORT_DATE")
                self._rename_shape(slide, "Text 9", "TXT_ASM_NAME")
                self._rename_shape(slide, "Text 11", "TXT_CHT_NAME")
                
            elif slide_id == "STORE_GENERAL_INFO":
                self._rename_shape(slide, "Text 10", "TXT_STORE_NAME")
                self._rename_shape(slide, "Text 13", "TXT_STORE_ADDRESS")
                self._rename_shape(slide, "Text 16", "TXT_REPORT_DATE")
                self._rename_shape(slide, "Text 19", "TXT_TIME_START")
                self._rename_shape(slide, "Text 22", "TXT_TIME_END")
                self._rename_shape(slide, "Text 25", "TXT_ASM_NAME")
                self._rename_shape(slide, "Text 28", "TXT_CHT_NAME")
                self._rename_shape(slide, "Text 31", "TXT_CHP_NAME")
                self._rename_shape(slide, "Shape 37", "TXT_NV_1")
                self._rename_shape(slide, "Shape 39", "TXT_NV_2")
                self._rename_shape(slide, "Shape 41", "TXT_NV_3")
                self._rename_shape(slide, "Shape 43", "TXT_NV_4")
                self._rename_shape(slide, "Shape 45", "TXT_BV_1")
                self._rename_shape(slide, "Shape 47", "TXT_BV_2")
                self._rename_shape(slide, "Text 50", "TXT_GENERAL_COMMENT")
                
            elif slide_id == "STORE_FRONTAGE_PHOTOS":
                self._rename_shape(slide, "Text 39", "TXT_FRONTAGE_COMMENTS")
                self._rename_shape(slide, "Shape 5", "PIC_FRONTAGE_1")
                self._rename_shape(slide, "Text 6", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_1")
                self._rename_shape(slide, "Shape 8", "PIC_FRONTAGE_2")
                self._rename_shape(slide, "Text 9", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_2")
                self._rename_shape(slide, "Shape 11", "PIC_FRONTAGE_3")
                self._rename_shape(slide, "Text 12", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_3")
                self._rename_shape(slide, "Shape 31", "SHP_RATING_TỐT")
                self._rename_shape(slide, "Shape 33", "SHP_RATING_ĐẠT")
                self._rename_shape(slide, "Shape 35", "SHP_RATING_CHƯAĐẠT")
                
            elif slide_id == "STORE_INNER_PHOTOS":
                self._rename_shape(slide, "Shape 5", "PIC_INNER_1")
                self._rename_shape(slide, "Text 6", "TXT_INNER_IMAGE_PLACEHOLDER_1")
                self._rename_shape(slide, "Shape 8", "PIC_INNER_2")
                self._rename_shape(slide, "Text 9", "TXT_INNER_IMAGE_PLACEHOLDER_2")
                self._rename_shape(slide, "Shape 11", "PIC_INNER_3")
                self._rename_shape(slide, "Text 12", "TXT_INNER_IMAGE_PLACEHOLDER_3")
                self._rename_shape(slide, "Text 39", "TXT_INNER_COMMENTS")
                self._rename_shape(slide, "Shape 31", "SHP_RATING_TỐT")
                self._rename_shape(slide, "Shape 33", "SHP_RATING_ĐẠT")
                self._rename_shape(slide, "Shape 35", "SHP_RATING_CHƯAĐẠT")
                
            elif slide_id in ["STORE_VM_ERROR_1", "STORE_VM_ERROR_2", "STORE_VM_ERROR_3", "STORE_VM_ERROR_4"]:
                self._rename_shape(slide, "Shape 7", "PIC_VM_BEFORE")
                self._rename_shape(slide, "Text 8", "TXT_VM_IMAGE_PLACEHOLDER_BEFORE")
                self._rename_shape(slide, "Shape 11", "PIC_VM_AFTER")
                self._rename_shape(slide, "Text 12", "TXT_VM_IMAGE_PLACEHOLDER_AFTER")
                self._rename_shape(slide, "Shape 15", "PIC_VM_DETAIL")
                self._rename_shape(slide, "Text 16", "TXT_VM_IMAGE_PLACEHOLDER_DETAIL")
                self._rename_shape(slide, "Text 41", "TXT_VM_ERROR_COMMENT")
                self._rename_shape(slide, "Shape 34", "SHP_RATING_TỐT")
                self._rename_shape(slide, "Shape 36", "SHP_RATING_ĐẠT")
                self._rename_shape(slide, "Shape 38", "SHP_RATING_CHƯAĐẠT")
                
            elif slide_id == "STORE_STOCKROOM_CASHIER":
                self._rename_shape(slide, "Shape 8", "PIC_STOCKROOM")
                self._rename_shape(slide, "Text 9", "TXT_STOCKROOM_IMAGE_PLACEHOLDER")
                self._rename_shape(slide, "Shape 11", "PIC_FITTING_ROOM")
                self._rename_shape(slide, "Text 12", "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER")
                self._rename_shape(slide, "Shape 29", "PIC_CASHIER")
                self._rename_shape(slide, "Text 30", "TXT_CASHIER_IMAGE_PLACEHOLDER")
                
            elif slide_id == "STORE_COMPETITOR":
                self._rename_shape(slide, "Shape 5", "PIC_COMP_1")
                self._rename_shape(slide, "Text 6", "TXT_COMP_IMAGE_PLACEHOLDER_1")
                self._rename_shape(slide, "Shape 8", "PIC_COMP_2")
                self._rename_shape(slide, "Text 9", "TXT_COMP_IMAGE_PLACEHOLDER_2")
                self._rename_shape(slide, "Shape 11", "PIC_COMP_3")
                self._rename_shape(slide, "Text 12", "TXT_COMP_IMAGE_PLACEHOLDER_3")
                self._rename_shape(slide, "Text 18", "TXT_COMP_TRAFFIC")
                self._rename_shape(slide, "Text 20", "TXT_COMP_COMPARISON")
                self._rename_shape(slide, "Text 22", "TXT_COMP_PEAK_TIME")
                self._rename_shape(slide, "Text 24", "TXT_COMP_ATTRACTION")
                self._rename_shape(slide, "Text 26", "TXT_COMP_SERVICE_APPEARANCE")
                self._rename_shape(slide, "Text 39", "TXT_COMP_ANALYSIS_SOLUTION")
                
            elif slide_id == "STORE_REVENUE":
                self._rename_shape(slide, "Text 8", "KPI_REVENUE_ACTUAL")
                self._rename_shape(slide, "Text 12", "KPI_REVENUE_TARGET")
                self._rename_shape(slide, "Text 16", "KPI_REVENUE_ATTAINMENT")
                self._rename_shape(slide, "Text 20", "KPI_REVENUE_PREV")
                self._rename_shape(slide, "Text 24", "KPI_REVENUE_YOY")
                self._rename_shape(slide, "Text 29", "TXT_REVENUE_COMMENT")
                
            elif slide_id == "STORE_STOCK_INVENTORY":
                self._rename_shape(slide, "Text 8", "KPI_STOCK_TOTAL")
                self._rename_shape(slide, "Text 16", "KPI_STOCK_NGUYEN_GIA")
                self._rename_shape(slide, "Text 20", "KPI_STOCK_SALE")
                self._rename_shape(slide, "Text 24", "KPI_STOCK_THANH_LY")
                # Do NOT rename TextBox 1 to TXT_STOCK_COMMENT to avoid overwriting aging groups text.
                
            elif slide_id == "STORE_BEST_SELLERS":
                for s in slide.shapes:
                    if s.has_text_frame and "TOP 10 SẢN PHẨM BÁN CHẠY" in s.text_frame.text:
                        s.name = "TBL_BEST_SELLERS"
                        break
                        
            elif slide_id == "STORE_SLOW_SELLERS":
                for s in slide.shapes:
                    if s.has_text_frame and "TOP 10 SẢN PHẨM BÁN CHẬM" in s.text_frame.text:
                        s.name = "TBL_SLOW_SELLERS"
                        break
                        
            elif slide_id == "STORE_STAFF_LIST":
                for s in slide.shapes:
                    if s.has_table:
                        s.name = "TBL_STORE_STAFF"
                        break
                        
            elif slide_id == "STORE_PENDING_ISSUES":
                for s in slide.shapes:
                    if s.has_table:
                        s.name = "TBL_PENDING_ISSUES"
                        break
                self._rename_shape(slide, "Text 115", "TXT_CSVC_ISSUE_NOTE")
                self._rename_shape(slide, "Shape 116", "PIC_CSVC_ISSUE")
                self._rename_shape(slide, "Text 117", "TXT_CSVC_IMAGE_PLACEHOLDER")
                self._rename_shape(slide, "Text 20", "TXT_ACTION_PLAN_PROPOSALS")
                
            elif slide_id == "STORE_DEV_PROPOSALS":
                self._rename_shape(slide, "Text 12", "TXT_DEV_PROPOSALS")
                
            elif slide_id == "STORE_THANK_YOU":
                self._rename_shape(slide, "Text 3", "TXT_FOOTER_INFO")

        prs.save(store_path)
        print(f"Initialized Store master template shapes and saved to {store_path}")

        # 2. Setup Cluster Master
        cluster_path = os.path.join(self.templates_dir, self.manifest["templates"]["cluster_master"])
        if not os.path.exists(cluster_path):
            raise FileNotFoundError(f"Cluster master template not found at {cluster_path}")
            
        prs_c = Presentation(cluster_path)
        
        if len(prs_c.slides) > 0:
            self._inject_metadata_shape(prs_c.slides[0], "CLUSTER_COVER")
            
        if len(prs_c.slides) > 1:
            self._inject_metadata_shape(prs_c.slides[1], "CLUSTER_SUMMARY")
                
        prs_c.save(cluster_path)
        print(f"Initialized Cluster master template shapes and saved to {cluster_path}")

    def verify_templates(self):
        """
        Verify template version, slide IDs, shape name uniqueness, and required shapes.
        """
        # Verify Store Master
        store_path = os.path.join(self.templates_dir, self.manifest["templates"]["store_master"])
        self._verify_single_template(store_path, "store_master")
        
        # Verify Cluster Master
        cluster_path = os.path.join(self.templates_dir, self.manifest["templates"]["cluster_master"])
        self._verify_single_template(cluster_path, "cluster_master")

    def _verify_single_template(self, path: str, type_label: str):
        if not os.path.exists(path):
            raise PreflightError(f"Template missing: {path}")
            
        try:
            prs = Presentation(path)
        except Exception as e:
            raise PreflightError(f"Corrupted PowerPoint file {path}: {str(e)}")
            
        # Discover slide metadata and verify shape uniqueness
        found_slide_ids = set()
        for i, slide in enumerate(prs.slides):
            meta_id = self._get_slide_metadata_id(slide)
            if not meta_id:
                continue
                
            # Uniqueness of Slide IDs
            if meta_id in found_slide_ids:
                raise PreflightError(f"Duplicate slide metadata ID '{meta_id}' found in {path}")
            found_slide_ids.add(meta_id)
            
            shape_names = [s.name for s in slide.shapes if s.name]
            
            # Check shape name uniqueness for technical names only (TXT_, PIC_, TBL_, KPI_, META_, SHP_)
            technical_names = [n for n in shape_names if n and any(n.startswith(prefix) for prefix in ("TXT_", "PIC_", "TBL_", "KPI_", "META_", "SHP_"))]
            if len(technical_names) != len(set(technical_names)):
                duplicates = set([n for n in technical_names if technical_names.count(n) > 1])
                raise PreflightError(f"Slide {i} (ID: {meta_id}) contains duplicate technical shape names: {duplicates}")
                
            # Verify required shapes are present
            manifest_slide = self.manifest["slides"].get(meta_id)
            if manifest_slide:
                required = manifest_slide.get("required_shapes", [])
                missing = [req for req in required if req not in shape_names]
                if missing:
                    raise PreflightError(f"Slide {i} (ID: {meta_id}) is missing required shapes: {missing}")

    def _get_slide_metadata_id(self, slide) -> str:
        for shape in slide.shapes:
            if shape.name == "META_SLIDE_ID" and shape.has_text_frame:
                return shape.text_frame.text.strip()
        return ""

    def _inject_metadata_shape(self, slide, slide_id: str):
        """Helper to inject hidden text frames with slide ID and version metadata. Fully idempotent."""
        id_shapes = [s for s in slide.shapes if s.name == "META_SLIDE_ID"]
        ver_shapes = [s for s in slide.shapes if s.name == "META_TEMPLATE_VERSION"]
        
        if len(id_shapes) > 1:
            raise PreflightError("Duplicate META_SLIDE_ID shapes found on slide.")
        if len(ver_shapes) > 1:
            raise PreflightError("Duplicate META_TEMPLATE_VERSION shapes found on slide.")
            
        # Update or create slide ID
        if len(id_shapes) == 1:
            id_shapes[0].text_frame.text = slide_id
        else:
            left = Inches(-10)
            top = Inches(-10)
            width = Inches(2)
            height = Inches(1)
            tx_box = slide.shapes.add_textbox(left, top, width, height)
            tx_box.name = "META_SLIDE_ID"
            tx_box.text_frame.text = slide_id
            
        # Update or create template version
        if len(ver_shapes) == 1:
            ver_shapes[0].text_frame.text = self.manifest["template_version"]
        else:
            left = Inches(-10)
            top = Inches(-10)
            width = Inches(2)
            height = Inches(1)
            tx_ver = slide.shapes.add_textbox(left, top + Inches(1.5), width, height)
            tx_ver.name = "META_TEMPLATE_VERSION"
            tx_ver.text_frame.text = self.manifest["template_version"]

    def _rename_shape(self, slide, old_name: str, new_name: str):
        """Find shape by old name or text and rename it."""
        for shape in slide.shapes:
            if shape.name == old_name:
                shape.name = new_name
                return
        # Fallback to matching by text content if old name not found
        for shape in slide.shapes:
            if shape.has_text_frame and old_name in shape.text_frame.text:
                shape.name = new_name
                return
