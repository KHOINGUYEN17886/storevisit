import os
import sys
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from data.models import StoreReportData
from .image_processor import ImageProcessor

# ── DESIGN TOKENS (SSOT màu & font cho toàn bộ báo cáo) ─────────────────────────
# Font chuẩn: Be Vietnam Pro (đã cài per-user + nên embed vào template master).
FONT_PRIMARY = "Be Vietnam Pro"
# Bảng màu thương hiệu An Phước — dùng tường minh, KHÔNG kế thừa màu placeholder xám.
CLR_NAVY   = RGBColor(0x0A, 0x23, 0x42)   # tiêu đề / heading
CLR_INK    = RGBColor(0x1A, 0x1A, 0x1A)   # body text (gần đen, rõ nét)
CLR_MUTED  = RGBColor(0x55, 0x5F, 0x6B)   # chú thích phụ (đủ đậm, không mờ)
CLR_OK     = RGBColor(0x1E, 0x8E, 0x3E)   # Tốt / Đạt
CLR_WARN   = RGBColor(0xF2, 0x99, 0x00)   # cảnh báo
CLR_ERR    = RGBColor(0xC0, 0x39, 0x2B)   # Chưa đạt / lỗi


def _is_placeholder_gray(rgb) -> bool:
    """True nếu màu là xám nhạt kiểu hint placeholder (cần thay bằng màu đậm).
    Giữ nguyên các màu bão hoà (KPI, rating) và màu tối thực sự."""
    try:
        r, g, b = rgb[0], rgb[1], rgb[2]
    except Exception:
        return False
    mx, mn = max(r, g, b), min(r, g, b)
    is_grayish = (mx - mn) <= 28          # gần như xám (R≈G≈B)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return is_grayish and luminance >= 110  # xám sáng → coi là placeholder hint


class PPTXGenerator:
    def __init__(self, template_path: str):
        self.template_path = template_path
        # Narrator LLM (Gemini→Groq→Claude→rule-based); tự tắt nếu không có API key.
        try:
            from reports.narrator import get_narrator
            self.narrator = get_narrator()
        except Exception as _e:
            self.narrator = None
            print(f"[PPTXGenerator] Narrator không khả dụng, dùng template: {_e}")

    def generate_store_report(self, data: StoreReportData, output_path: str, temp_image_paths: dict):
        """
        Populate individual store slides in the presentation.
        Finds slides using META_SLIDE_ID rather than hardcoded slide indexes.
        """
        prs = Presentation(self.template_path)
        
        # Build map of slide ID to Slide object
        slide_map = {}
        for slide in prs.slides:
            slide_id = self._get_slide_id(slide)
            if slide_id:
                slide_map[slide_id] = slide

        slides_to_delete = []
        
        # Parse checklist_json
        c_data = {}
        brand_presence = {}
        if data.form_response and data.form_response.checklist_json:
            try:
                import json
                c_data = json.loads(data.form_response.checklist_json)
                brand_presence = c_data.get("brand_presence", {})
            except Exception as e:
                print(f"Error parsing checklist_json: {e}")

        # Determine if data is missing
        missing_revenue = temp_image_paths.get("missing_revenue", False)
        missing_stock = temp_image_paths.get("missing_stock", False)

        # Đồng bộ 30-07: phân loại lượt kiểm tra (own/cross/opening) — dùng chung cho
        # nhãn trên slide thông tin chung VÀ để quyết định ẩn slide doanh thu/tồn kho
        # khi là báo cáo khai trương (xem bên dưới, gần STORE_REVENUE).
        inspection_mode = ((data.form_response.inspection_mode if data.form_response else None) or "own").strip().lower()
        is_opening_report = (inspection_mode == "opening")

        # 1. STORE_COVER
        cover_slide = slide_map.get("STORE_COVER")
        if cover_slide:
            self._fill_text(cover_slide, "TXT_STORE_NAME", f"Cửa hàng An Phước {data.metadata.store_name}", font_size=20, bold=True, align="center")
            report_date = data.form_response.report_date if data.form_response and data.form_response.report_date else "-"
            self._fill_text(cover_slide, "TXT_REPORT_DATE", report_date, font_size=16, align="center")
            asm_name = data.form_response.asm_name if data.form_response and data.form_response.asm_name else data.metadata.asm_name
            self._fill_text(cover_slide, "TXT_ASM_NAME", asm_name, font_size=16, align="center")
            cht_name = data.form_response.cht_name if data.form_response and data.form_response.cht_name else data.staff.cht_name
            self._fill_text(cover_slide, "TXT_CHT_NAME", cht_name, font_size=16, align="center")

        # 2. STORE_GENERAL_INFO
        info_slide = slide_map.get("STORE_GENERAL_INFO")
        if info_slide:
            self._fill_text(info_slide, "TXT_STORE_NAME", f"An Phước {data.metadata.store_name}")
            self._fill_text(info_slide, "TXT_STORE_ADDRESS", data.metadata.address or "-")
            
            report_date = data.form_response.report_date if data.form_response and data.form_response.report_date else "-"
            time_start = data.form_response.time_start if data.form_response and data.form_response.time_start else "-"
            time_end = data.form_response.time_end if data.form_response and data.form_response.time_end else "-"
            asm_name = data.form_response.asm_name if data.form_response and data.form_response.asm_name else data.metadata.asm_name
            cht_name = data.form_response.cht_name if data.form_response and data.form_response.cht_name else data.staff.cht_name
            
            self._fill_text(info_slide, "TXT_REPORT_DATE", report_date)
            self._fill_text(info_slide, "TXT_TIME_START", time_start)
            self._fill_text(info_slide, "TXT_TIME_END", time_end)
            self._fill_text(info_slide, "TXT_ASM_NAME", asm_name)
            self._fill_text(info_slide, "TXT_CHT_NAME", cht_name)
            self._fill_text(info_slide, "TXT_CHP_NAME", data.staff.chp_name or "-")
            
            # Staff names
            nvbh_list = [s for s in data.staff.staff_list if "bảo vệ" not in s.role.lower() and s.name not in [cht_name, data.staff.chp_name]]
            bv_list = [s for s in data.staff.staff_list if "bảo vệ" in s.role.lower()]
            
            # NVBH roster slots
            for i in range(4):
                shape_name = f"TXT_NV_{i+1}"
                if i < len(nvbh_list):
                    staff = nvbh_list[i]
                    self._fill_text(info_slide, shape_name, f"{staff.name} ({staff.role})")
                else:
                    self._fill_text(info_slide, shape_name, "(Trống)")
                    
            # BV slots
            for i in range(2):
                shape_name = f"TXT_BV_{i+1}"
                if i < len(bv_list):
                    self._fill_text(info_slide, shape_name, bv_list[i].name)
                else:
                    self._fill_text(info_slide, shape_name, "(Trống)")

            if is_opening_report:
                # Báo cáo khai trương: thay hẳn "nhận định chung" (vốn dựa trên doanh thu/tồn
                # kho — không có ý nghĩa với cửa hàng mới mở/đang sửa chữa) bằng thông tin
                # khai trương + badge mức độ sẵn sàng màu theo CLR_OK/CLR_WARN/CLR_ERR.
                self._fill_opening_info(info_slide, data.form_response)
            else:
                # Overall comment: "Nhận định chung của QLKD" — humanized qua narrator, grounded theo dữ liệu.
                _general_fallback = "Không ghi nhận nhận xét tổng quan bổ sung tại thời điểm kiểm tra."
                general_text = _general_fallback
                try:
                    facts = self._build_exec_facts(data, c_data)
                    if facts and self.narrator is not None:
                        general_text = self.narrator.executive_summary(
                            facts=facts,
                            store_name=data.metadata.store_name,
                            fallback=_general_fallback,
                        )
                except Exception as _e:
                    print(f"[narrator] exec summary lỗi, dùng fallback: {_e}")
                if inspection_mode == "cross":
                    general_text = f"🔀 KIỂM TRA CHÉO — thực hiện bởi {asm_name} (thay ASM phụ trách cửa hàng này).\n\n{general_text}"
                self._fill_text(info_slide, "TXT_GENERAL_COMMENT", general_text)

        # 3. STORE_FRONTAGE_PHOTOS & STORE_INNER_PHOTOS (Merging if both have no photos)
        frontage_slide = slide_map.get("STORE_FRONTAGE_PHOTOS")
        inner_slide = slide_map.get("STORE_INNER_PHOTOS")
        
        # Determine frontage comments & photos
        front_photos = []
        front_rating = data.frontage_rating or "Đạt"
        front_comment = "Không ghi nhận vấn đề mặt tiền cần khắc phục tại thời điểm kiểm tra."
        if frontage_slide:
            if data.form_response:
                front_photos = sorted([p for p in data.form_response.photos if p.section == "frontage"], key=lambda x: x.index)
                if data.form_response.rating_frontage:
                    front_rating = data.form_response.rating_frontage
                
                general_comment = data.form_response.comment_frontage or ""
                resolved_info = self._collect_resolved_comments(c_data, ["frontage"])
                parts = []
                if general_comment:
                    parts.append(f"Nhận xét chung: {general_comment}")
                if resolved_info["failures"]:
                    parts.append("Lỗi chưa khắc phục:\n" + "\n".join(resolved_info["failures"]))
                if resolved_info["resolved"]:
                    parts.append("Lỗi đã khắc phục tại chỗ:\n" + "\n".join(resolved_info["resolved"]))
                if parts:
                    front_comment = "\n\n".join(parts)
            else:
                if data.frontage_issue:
                    if data.frontage_action:
                        front_comment = f"Lỗi phát hiện: {data.frontage_issue}\n-> Hướng khắc phục: {data.frontage_action}"
                    else:
                        front_comment = data.frontage_issue

        # Determine inner comments & photos
        inner_photos = []
        inner_comment = "Không ghi nhận lỗi trưng bày hoặc không gian bên trong."
        inner_rating = "Đạt"
        if inner_slide:
            if data.form_response:
                inner_photos = sorted([p for p in data.form_response.photos if p.section == "inner"], key=lambda x: x.index)
                if data.form_response.rating_inner:
                    inner_rating = data.form_response.rating_inner
                
                general_comment = data.form_response.comment_inner or data.form_response.comment_merch or ""
                resolved_info = self._collect_resolved_comments(c_data, ["inner"])
                parts = []
                if general_comment:
                    parts.append(f"Nhận xét chung: {general_comment}")
                if resolved_info["failures"]:
                    parts.append("Lỗi chưa khắc phục:\n" + "\n".join(resolved_info["failures"]))
                if resolved_info["resolved"]:
                    parts.append("Lỗi đã khắc phục tại chỗ:\n" + "\n".join(resolved_info["resolved"]))
                if parts:
                    inner_comment = "\n\n".join(parts)

        # Check physical photo existence
        has_front_photo = any(p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0 for p in front_photos)
        has_inner_photo = any(p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0 for p in inner_photos)

        # Đồng bộ 30-07: báo cáo Tái khai trương cần GIỮ inner_slide sống để ghi đè bằng
        # so sánh ảnh Trước/Sau bên dưới — không cho merge-xoá dù frontage/inner rỗng ảnh.
        is_reopen_report = is_opening_report and data.form_response and data.form_response.opening_type == "reopen"
        if frontage_slide and inner_slide and not has_front_photo and not has_inner_photo and not is_reopen_report:
            # Merge! Delete existing body shapes on frontage slide
            shapes_to_delete = []
            for s in frontage_slide.shapes:
                if s.name not in ["META_SLIDE_ID", "META_TEMPLATE_VERSION"] and s.top and s.top > Inches(1.85):
                    if s.has_text_frame and ("Cửa hàng:" in s.text_frame.text or "Ngày:" in s.text_frame.text):
                        continue
                    shapes_to_delete.append(s)
            for s in shapes_to_delete:
                try: frontage_slide.shapes._spTree.remove(s._element)
                except Exception: pass
                
            # Create a clean two-column layout
            # Column 1: Frontage
            tb_1 = frontage_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.5), Inches(7.5))
            tf_1 = tb_1.text_frame
            tf_1.word_wrap = True
            
            p1 = tf_1.paragraphs[0]
            p1.text = "I. ĐÁNH GIÁ MẶT TIỀN CỬA HÀNG"
            p1.font.name = FONT_PRIMARY
            p1.font.size = Pt(14)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(12, 35, 64)
            p1.space_after = Pt(12)
            
            p1_rate = tf_1.add_paragraph()
            p1_rate.text = f"Xếp loại ngoại quan: {front_rating}"
            p1_rate.font.name = FONT_PRIMARY
            p1_rate.font.size = Pt(11)
            p1_rate.font.bold = True
            p1_rate.font.color.rgb = RGBColor(220, 100, 50) if front_rating == "Chưa đạt" else RGBColor(40, 160, 80)
            p1_rate.space_after = Pt(12)
            
            self._add_multiline_paragraphs(tf_1, front_comment, font_name=FONT_PRIMARY, font_size=10.5, color_rgb=RGBColor(0, 0, 0))
            
            # Column 2: Inner
            tb_2 = frontage_slide.shapes.add_textbox(Inches(10.0), Inches(2.2), Inches(8.5), Inches(7.5))
            tf_2 = tb_2.text_frame
            tf_2.word_wrap = True
            
            p2 = tf_2.paragraphs[0]
            p2.text = "II. ĐÁNH GIÁ KHÔNG GIAN BÊN TRONG"
            p2.font.name = FONT_PRIMARY
            p2.font.size = Pt(14)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(12, 35, 64)
            p2.space_after = Pt(12)
            
            p2_rate = tf_2.add_paragraph()
            p2_rate.text = f"Xếp loại bên trong: {inner_rating}"
            p2_rate.font.name = FONT_PRIMARY
            p2_rate.font.size = Pt(11)
            p2_rate.font.bold = True
            p2_rate.font.color.rgb = RGBColor(220, 100, 50) if inner_rating == "Chưa đạt" else RGBColor(40, 160, 80)
            p2_rate.space_after = Pt(12)
            
            self._add_multiline_paragraphs(tf_2, inner_comment, font_name=FONT_PRIMARY, font_size=10.5, color_rgb=RGBColor(0, 0, 0))
            
            # Rename slide main header title
            for s in frontage_slide.shapes:
                if s.has_text_frame and ("MẶT TIỀN" in s.text_frame.text or "NGOẠI QUAN" in s.text_frame.text):
                    s.text_frame.text = "2.1. ĐÁNH GIÁ MẶT TIỀN & BÊN TRONG CỬA HÀNG"
                    if s.text_frame.paragraphs and s.text_frame.paragraphs[0].runs:
                        s.text_frame.paragraphs[0].runs[0].font.name = FONT_PRIMARY
                        s.text_frame.paragraphs[0].runs[0].font.bold = True
            
            # Delete inner slide
            slides_to_delete.append(inner_slide)
        else:
            # Normal render (separate slides)
            if frontage_slide:
                self._highlight_rating(frontage_slide, front_rating)
                self._fill_text(frontage_slide, "TXT_FRONTAGE_COMMENTS", front_comment)
                slot_configs = [
                    ("PIC_FRONTAGE_1", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_1"),
                    ("PIC_FRONTAGE_2", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_2"),
                    ("PIC_FRONTAGE_3", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_3"),
                ]
                photo_paths = [p.local_path for p in front_photos if p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]
                if len(photo_paths) <= 3:
                    padded_photo_paths = photo_paths + [""] * (3 - len(photo_paths))
                    self._fill_image_slots_adaptive(frontage_slide, slot_configs, padded_photo_paths)
                else:
                    self._fill_image_slots_grid(frontage_slide, slot_configs, photo_paths)
                
            if inner_slide:
                self._highlight_rating(inner_slide, inner_rating)
                self._fill_text(inner_slide, "TXT_INNER_COMMENTS", inner_comment)
                slot_configs = [
                    ("PIC_INNER_1", "TXT_INNER_IMAGE_PLACEHOLDER_1"),
                    ("PIC_INNER_2", "TXT_INNER_IMAGE_PLACEHOLDER_2"),
                    ("PIC_INNER_3", "TXT_INNER_IMAGE_PLACEHOLDER_3"),
                ]
                photo_paths = [p.local_path for p in inner_photos if p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]
                if len(photo_paths) <= 3:
                    padded_photo_paths = photo_paths + [""] * (3 - len(photo_paths))
                    self._fill_image_slots_adaptive(inner_slide, slot_configs, padded_photo_paths)
                else:
                    self._fill_image_slots_grid(inner_slide, slot_configs, photo_paths)

        # Đồng bộ 30-07: báo cáo Tái khai trương — ghi đè slide STORE_INNER_PHOTOS (vừa
        # điền ở trên) bằng so sánh ảnh Trước/Sau sửa chữa. Tái dùng slide có sẵn trong
        # template thay vì cần slide mới (python-pptx không dễ chèn slide layout mới lúc chạy).
        if is_reopen_report and inner_slide:
            self._build_opening_before_after(inner_slide, data.form_response)

        # 5. STORE_VM_ERROR_1 to 4
        vm_mapping = {
            0: ("vm_ap", "An Phước", "merch_ap"),
            1: ("vm_pie", "Pierre Cardin", "merch_pie"),
            2: ("vm_ab", "Anamai/Bonjour", "merch_ab"),
            3: ("vm_pk", "Phụ kiện", "merch_pk")
        }
        for idx in range(4):
            vm_id = f"STORE_VM_ERROR_{idx+1}"
            vm_slide = slide_map.get(vm_id)
            if vm_slide:
                sec_name, label, json_key = vm_mapping[idx]
                
                # Check brand presence
                if idx == 0:
                    is_present = brand_presence.get("merch_ap", True)
                elif idx == 1:
                    is_present = brand_presence.get("merch_pie", True)
                elif idx == 2:
                    is_present = brand_presence.get("merch_anamai", True) or brand_presence.get("merch_bonjour", True)
                else:
                    is_present = brand_presence.get("merch_pk", True)
                    
                # Determine VM rating
                vm_rating = "Đạt"
                sec_val = {}
                if c_data:
                    if json_key == "merch_ab" and "merch_ab" not in c_data.get("sections", {}):
                        anamai_sec = c_data.get("sections", {}).get("merch_anamai", {})
                        bonjour_sec = c_data.get("sections", {}).get("merch_bonjour", {})
                        r1 = anamai_sec.get("rating", "Đạt")
                        r2 = bonjour_sec.get("rating", "Đạt")
                        if "Chưa đạt" in [r1, r2]:
                            vm_rating = "Chưa đạt"
                        elif r1 == "Tốt" and r2 == "Tốt":
                            vm_rating = "Tốt"
                        else:
                            vm_rating = "Đạt"
                        combined_items = anamai_sec.get("items", []) + bonjour_sec.get("items", [])
                        sec_val = {"rating": vm_rating, "items": combined_items}
                    else:
                        sec_val = c_data.get("sections", {}).get(json_key, {})
                        vm_rating = sec_val.get("rating", "Đạt")
                        
                # Delete slide if not present or rating is not failure (meaning Đạt, Tốt, Không áp dụng)
                if not is_present or vm_rating in ["Đạt", "Tốt", "Không áp dụng"]:
                    slides_to_delete.append(vm_slide)
                    continue
                
                vm_photos = []
                if data.form_response:
                    vm_photos = [p for p in data.form_response.photos if p.section == sec_name]
                    
                slot_configs_vm = [
                    ("PIC_VM_BEFORE", "TXT_VM_IMAGE_PLACEHOLDER_BEFORE"),
                    ("PIC_VM_AFTER",  "TXT_VM_IMAGE_PLACEHOLDER_AFTER"),
                    ("PIC_VM_DETAIL", "TXT_VM_IMAGE_PLACEHOLDER_DETAIL"),
                ]
                photo_paths_vm = [p.local_path for p in vm_photos if p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]
                if len(photo_paths_vm) <= 3:
                    padded_photo_paths_vm = photo_paths_vm + [""] * (3 - len(photo_paths_vm))
                    self._fill_image_slots_adaptive(vm_slide, slot_configs_vm, padded_photo_paths_vm)
                else:
                    self._fill_image_slots_grid(vm_slide, slot_configs_vm, photo_paths_vm)
                
                comment = f"Lỗi phát hiện: Không ghi nhận lỗi trưng bày tại khu vực {label}.\n→ Hướng khắc phục: Không phát sinh."
                if c_data:
                    try:
                        failures = []
                        resolutions = []
                        for item in sec_val.get("items", []):
                            if item.get("eval") == "Không đạt":
                                item_label = item.get("label", "")
                                note = item.get("note", "Lỗi trưng bày")
                                severity = item.get("severity", "Trung bình")
                                resolved = item.get("resolved") or "Không"
                                assignee = item.get("assignee") or "CHT"
                                deadline = item.get("deadline") or "-"
                                
                                if resolved == "Có":
                                    failures.append(f"- {item_label}: {note} (Đã khắc phục tại chỗ)")
                                    resolutions.append(f"- {item_label}: Đã khắc phục tại chỗ.")
                                else:
                                    failures.append(f"- {item_label}: {note} ({severity})")
                                    resolutions.append(f"- {item_label}: CHT {assignee} chịu trách nhiệm, hạn {deadline}.")
                                    
                        if failures:
                            failures_str = "\n".join(failures)
                            resolutions_str = "\n".join(resolutions)
                            comment = f"Lỗi phát hiện:\n{failures_str}\n\n→ Hướng khắc phục:\n{resolutions_str}"
                    except Exception as e:
                        print(f"Error parsing checklist VM comment for {sec_name}: {e}")
                elif idx == 0 and data.form_response and data.form_response.comment_merch:
                    comment = data.form_response.comment_merch
                
                self._highlight_rating(vm_slide, vm_rating)
                self._fill_text(vm_slide, "TXT_VM_ERROR_COMMENT", comment)

        # Calculate operational comments and photo presence early
        has_stock_photo = False
        has_fit_photo = False
        has_cash_photo = False
        has_csvc_photo = False
        
        stockroom_photos = []
        fitting_photos = []
        cashier_photos = []
        csvc_photos = []
        if data.form_response:
            stockroom_photos = [p for p in data.form_response.photos if p.section == "stockroom" and p.local_path]
            fitting_photos = [p for p in data.form_response.photos if p.section == "fitting_room" and p.local_path]
            cashier_photos = [p for p in data.form_response.photos if p.section == "cashier" and p.local_path]
            csvc_photos = [p for p in data.form_response.photos if p.section == "csvc" and p.local_path]
            
        path_stock = stockroom_photos[0].local_path if stockroom_photos else ""
        path_fit = fitting_photos[0].local_path if fitting_photos else ""
        path_cash = cashier_photos[0].local_path if cashier_photos else ""
        csvc_path = csvc_photos[0].local_path if csvc_photos else ""
        
        has_stock_photo = bool(path_stock and os.path.exists(path_stock) and os.path.getsize(path_stock) > 0)
        has_fit_photo = bool(path_fit and os.path.exists(path_fit) and os.path.getsize(path_fit) > 0)
        has_cash_photo = bool(path_cash and os.path.exists(path_cash) and os.path.getsize(path_cash) > 0)
        has_csvc_photo = bool(csvc_path and os.path.exists(csvc_path) and os.path.getsize(csvc_path) > 0)
        
        warehouse_comment = "Lỗi khu vực kho phát hiện & Hướng xử lý: Không ghi nhận lỗi."
        cashier_comment = "Lỗi quầy thu ngân & bảo vệ phát hiện & Hướng xử lý: Không ghi nhận lỗi."
        
        if c_data:
            try:
                w_info = self._collect_resolved_comments(c_data, ["warehouse", "stockroom", "fitting_room", "toilet", "fire_safety"])
                w_parts = []
                if w_info["failures"]:
                    w_parts.append("Lỗi chưa khắc phục:\n" + "\n".join(w_info["failures"]))
                if w_info["resolved"]:
                    w_parts.append("Lỗi đã khắc phục tại chỗ:\n" + "\n".join(w_info["resolved"]))
                if w_parts:
                    warehouse_comment = "Lỗi khu vực kho phát hiện & Hướng xử lý:\n" + "\n\n".join(w_parts)
                    
                # Đồng bộ 30-07: "security_guard" gộp chung nhóm này — không có card riêng
                # trong layout 3-card cố định bên dưới, nhóm cùng cashier/packaging_security
                # (đều là hạng mục vận hành mặt tiền/khách hàng) để không bị rơi khỏi báo cáo.
                c_info = self._collect_resolved_comments(c_data, ["cashier", "packaging_security", "security_guard"])
                c_parts = []
                if c_info["failures"]:
                    c_parts.append("Lỗi chưa khắc phục:\n" + "\n".join(c_info["failures"]))
                if c_info["resolved"]:
                    c_parts.append("Lỗi đã khắc phục tại chỗ:\n" + "\n".join(c_info["resolved"]))
                if c_parts:
                    cashier_comment = "Lỗi quầy thu ngân & bảo vệ phát hiện & Hướng xử lý:\n" + "\n\n".join(c_parts)
            except Exception as e:
                print(f"Error parsing warehouse/cashier comments early: {e}")
                
        csvc_comment = data.csvc_comment or "Không ghi nhận vấn đề CSVC khẩn cấp cần chuyển PTTT."
        merge_operational_csvc = (not has_stock_photo and not has_fit_photo and not has_cash_photo and not has_csvc_photo)

        # 6. STORE_STOCKROOM_CASHIER
        stock_cashier_slide = slide_map.get("STORE_STOCKROOM_CASHIER")
        if stock_cashier_slide:
            if merge_operational_csvc:
                slides_to_delete.append(stock_cashier_slide)
            else:
                # Find label shapes
                label_stock = None
                label_fit = None
                label_cash = None
                for s in stock_cashier_slide.shapes:
                    if s.has_text_frame:
                        txt = s.text_frame.text.strip()
                        if txt == "Ảnh kho":
                            label_stock = s
                        elif txt == "Ảnh phòng thử":
                            label_fit = s
                        elif txt == "Ảnh quầy thu ngân":
                            label_cash = s

                # --- Process left half (Warehouse & Fitting room) ---
                valid_left = []
                if path_stock and os.path.exists(path_stock) and os.path.getsize(path_stock) > 0:
                    valid_left.append(("stock", path_stock))
                if path_fit and os.path.exists(path_fit) and os.path.getsize(path_fit) > 0:
                    valid_left.append(("fit", path_fit))
                    
                if len(valid_left) == 2:
                    self._fill_image_slot(stock_cashier_slide, "PIC_STOCKROOM", "TXT_STOCKROOM_IMAGE_PLACEHOLDER", path_stock)
                    self._fill_image_slot(stock_cashier_slide, "PIC_FITTING_ROOM", "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER", path_fit)
                elif len(valid_left) == 1:
                    # 1 photo: center it horizontally in the left column
                    type_img, img_path = valid_left[0]
                    new_width = Inches(5.5)
                    new_left = Inches(1.03 + (8.29 - 5.5) / 2) # = Inches(2.42)
                    
                    if type_img == "stock":
                        self._remove_shape_by_name(stock_cashier_slide, "PIC_FITTING_ROOM")
                        self._remove_shape_by_name(stock_cashier_slide, "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER")
                        if label_fit:
                            try: stock_cashier_slide.shapes._spTree.remove(label_fit._element)
                            except Exception: pass
                        
                        # Update PIC_STOCKROOM and label position
                        for s in stock_cashier_slide.shapes:
                            if s.name == "PIC_STOCKROOM":
                                s.left = new_left
                                s.width = new_width
                        if label_stock:
                            label_stock.left = new_left
                            label_stock.width = new_width
                            label_stock.text_frame.text = "Ảnh kho/phòng thử"
                            if label_stock.text_frame.paragraphs and label_stock.text_frame.paragraphs[0].runs:
                                label_stock.text_frame.paragraphs[0].runs[0].font.name = FONT_PRIMARY
                                
                        self._fill_image_slot(stock_cashier_slide, "PIC_STOCKROOM", "TXT_STOCKROOM_IMAGE_PLACEHOLDER", img_path)
                    else:
                        self._remove_shape_by_name(stock_cashier_slide, "PIC_STOCKROOM")
                        self._remove_shape_by_name(stock_cashier_slide, "TXT_STOCKROOM_IMAGE_PLACEHOLDER")
                        if label_stock:
                            try: stock_cashier_slide.shapes._spTree.remove(label_stock._element)
                            except Exception: pass
                            
                        # Update PIC_FITTING_ROOM and label position
                        for s in stock_cashier_slide.shapes:
                            if s.name == "PIC_FITTING_ROOM":
                                s.left = new_left
                                s.width = new_width
                        if label_fit:
                            label_fit.left = new_left
                            label_fit.width = new_width
                            
                        self._fill_image_slot(stock_cashier_slide, "PIC_FITTING_ROOM", "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER", img_path)
                else:
                    # 0 photos: hide all left photos and expand the warehouse comment box
                    self._remove_shape_by_name(stock_cashier_slide, "PIC_STOCKROOM")
                    self._remove_shape_by_name(stock_cashier_slide, "TXT_STOCKROOM_IMAGE_PLACEHOLDER")
                    self._remove_shape_by_name(stock_cashier_slide, "PIC_FITTING_ROOM")
                    self._remove_shape_by_name(stock_cashier_slide, "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER")
                    if label_stock:
                        try: stock_cashier_slide.shapes._spTree.remove(label_stock._element)
                        except Exception: pass
                    if label_fit:
                        try: stock_cashier_slide.shapes._spTree.remove(label_fit._element)
                        except Exception: pass
                        
                    # Expand warehouse comment box
                    kho_text_shape = None
                    for s in stock_cashier_slide.shapes:
                        if s.has_text_frame and s.text_frame.text.strip().startswith("Lỗi khu vực kho"):
                            kho_text_shape = s
                            break
                    if kho_text_shape:
                        kho_bg_shape = None
                        for s in stock_cashier_slide.shapes:
                            if s != kho_text_shape and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                if abs(s.top - kho_text_shape.top) < Inches(0.25) and abs(s.left - kho_text_shape.left) < Inches(0.25):
                                    kho_bg_shape = s
                                    break
                        
                        new_top = Inches(2.73)
                        new_height = Inches(9.98 - 2.73) # 7.25
                        if kho_bg_shape:
                            kho_bg_shape.top = new_top
                            kho_bg_shape.height = new_height
                        kho_text_shape.top = Inches(2.85)
                        kho_text_shape.height = Inches(6.80)
                        if kho_text_shape.has_text_frame:
                            kho_text_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP

                # --- Process right half (Cashier) ---
                has_cash = path_cash and os.path.exists(path_cash) and os.path.getsize(path_cash) > 0
                if has_cash:
                    self._fill_image_slot(stock_cashier_slide, "PIC_CASHIER", "TXT_CASHIER_IMAGE_PLACEHOLDER", path_cash)
                else:
                    self._remove_shape_by_name(stock_cashier_slide, "PIC_CASHIER")
                    self._remove_shape_by_name(stock_cashier_slide, "TXT_CASHIER_IMAGE_PLACEHOLDER")
                    if label_cash:
                        try: stock_cashier_slide.shapes._spTree.remove(label_cash._element)
                        except Exception: pass
                        
                    # Expand cashier comment box
                    tn_text_shape = None
                    for s in stock_cashier_slide.shapes:
                        if s.has_text_frame and s.text_frame.text.strip().startswith("Lỗi quầy thu ngân"):
                            tn_text_shape = s
                            break
                    if tn_text_shape:
                        tn_bg_shape = None
                        for s in stock_cashier_slide.shapes:
                            if s != tn_text_shape and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                if abs(s.top - tn_text_shape.top) < Inches(0.25) and abs(s.left - tn_text_shape.left) < Inches(0.25):
                                    tn_bg_shape = s
                                    break
                                    
                        new_top = Inches(2.73)
                        new_height = Inches(9.98 - 2.73)
                        if tn_bg_shape:
                            tn_bg_shape.top = new_top
                            tn_bg_shape.height = new_height
                        tn_text_shape.top = Inches(2.85)
                        tn_text_shape.height = Inches(6.80)
                        if tn_text_shape.has_text_frame:
                            tn_text_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                
                self._fill_text_by_prefix_or_name(stock_cashier_slide, "Lỗi khu vực kho", warehouse_comment, font_size=10)
                self._fill_text_by_prefix_or_name(stock_cashier_slide, "Lỗi quầy thu ngân", cashier_comment, font_size=10)

        # 7. STORE_COMPETITOR
        comp_slide = slide_map.get("STORE_COMPETITOR")
        if comp_slide:
            comp_photos = []
            if data.form_response:
                comp_photos = [p for p in data.form_response.photos if p.section == "competitor" and p.local_path]
                
            slot_configs = [
                ("PIC_COMP_1", "TXT_COMP_IMAGE_PLACEHOLDER_1"),
                ("PIC_COMP_2", "TXT_COMP_IMAGE_PLACEHOLDER_2"),
                ("PIC_COMP_3", "TXT_COMP_IMAGE_PLACEHOLDER_3"),
            ]
            photo_paths = [p.local_path for p in comp_photos if p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]
            if len(photo_paths) <= 3:
                padded_photo_paths = photo_paths + [""] * (3 - len(photo_paths))
                self._fill_image_slots_adaptive(comp_slide, slot_configs, padded_photo_paths)
            else:
                self._fill_image_slots_grid(comp_slide, slot_configs, photo_paths)

            comp_traffic = "Chưa ghi nhận thông tin tại thời điểm kiểm tra."
            comp_comparison = "Chưa ghi nhận thông tin tại thời điểm kiểm tra."
            comp_peak_time = "Chưa ghi nhận thông tin tại thời điểm kiểm tra."
            comp_attraction = "Chưa ghi nhận thông tin tại thời điểm kiểm tra."
            comp_service = "Chưa ghi nhận thông tin tại thời điểm kiểm tra."
            comp_comment = "Điểm mạnh/đặc điểm đối thủ: Chưa ghi nhận\n-> Giải pháp của QLKD: Chưa ghi nhận"
            
            if data.form_response and data.form_response.checklist_json:
                try:
                    import json
                    c_data = json.loads(data.form_response.checklist_json)
                    c_sec = c_data.get("competitor", {})
                    
                    if c_sec.get("traffic"):
                        comp_traffic = c_sec.get("traffic")
                    if c_sec.get("comparison"):
                        comp_comparison = c_sec.get("comparison")
                    if c_sec.get("peak_time"):
                        comp_peak_time = c_sec.get("peak_time")
                    if c_sec.get("attraction"):
                        comp_attraction = c_sec.get("attraction")
                    if c_sec.get("service_appearance"):
                        comp_service = c_sec.get("service_appearance")
                    if c_sec.get("analysis_solution"):
                        comp_comment = c_sec.get("analysis_solution")
                except Exception as e:
                    print(f"Error parsing competitor comments: {e}")
            elif hasattr(data, "weekly_json") and data.weekly_json:
                sec4 = data.weekly_json.get("sec4", [])
                comp_analysis = "Chưa ghi nhận"
                comp_solution = "Chưa ghi nhận"
                for item in sec4:
                    if item.get("label") and "đối thủ" in str(item["label"]).lower():
                        note = item.get("note") or item.get("issue")
                        if note:
                            comp_note = str(note).strip()
                            if "->" in comp_note:
                                comp_analysis, comp_solution = comp_note.split("->", 1)
                                comp_analysis = comp_analysis.strip()
                                comp_solution = comp_solution.strip()
                            else:
                                comp_analysis = comp_note
                comp_comment = f"Điểm mạnh/đặc điểm đối thủ: {comp_analysis}\n-> Giải pháp của QLKD: {comp_solution}"
                
            self._fill_text(comp_slide, "TXT_COMP_TRAFFIC", comp_traffic)
            self._fill_text(comp_slide, "TXT_COMP_COMPARISON", comp_comparison)
            self._fill_text(comp_slide, "TXT_COMP_PEAK_TIME", comp_peak_time)
            self._fill_text(comp_slide, "TXT_COMP_ATTRACTION", comp_attraction)
            self._fill_text(comp_slide, "TXT_COMP_SERVICE_APPEARANCE", comp_service)
            self._fill_text(comp_slide, "TXT_COMP_ANALYSIS_SOLUTION", comp_comment)

        # 8. STORE_REVENUE — ẩn với báo cáo khai trương (cửa hàng mới/đang sửa chữa
        # thường chưa có dữ liệu doanh thu ý nghĩa để so sánh, theo quyết định 30-07).
        rev_slide = slide_map.get("STORE_REVENUE")
        if rev_slide and is_opening_report:
            slides_to_delete.append(rev_slide)
        elif rev_slide:
            if missing_revenue:
                self._fill_text(rev_slide, "KPI_REVENUE_ACTUAL", "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_TARGET", "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_ATTAINMENT", "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_PREV", "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_YOY", "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "TXT_REVENUE_COMMENT", "Chưa có dữ liệu doanh thu trong kỳ báo cáo.")
            else:
                self._fill_text(rev_slide, "KPI_REVENUE_ACTUAL", f"{data.revenue.revenue_actual:,.0f} VNĐ" if data.revenue.revenue_actual is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_TARGET", f"{data.revenue.revenue_target:,.0f} VNĐ" if data.revenue.revenue_target is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_ATTAINMENT", f"{data.revenue.attainment_pct:.1f}%" if data.revenue.attainment_pct is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_PREV", f"{data.revenue.revenue_prev:,.0f} VNĐ" if data.revenue.revenue_prev is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "KPI_REVENUE_YOY", f"{data.revenue.revenue_yoy:,.0f} VNĐ" if data.revenue.revenue_yoy is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(rev_slide, "TXT_REVENUE_COMMENT", data.revenue.commentary or "-")
                
                # Set Action Title
                attainment = data.revenue.attainment_pct if data.revenue.attainment_pct is not None else 0.0
                self._set_slide_title(rev_slide, "3.1.", f"3.1. DOANH THU ĐẠT {attainment:.1f}% KẾ HOẠCH")
            
            # If a revenue chart path is provided, we can insert it over the OLE chart area
            chart_img = temp_image_paths.get("revenue_chart", "")
            if chart_img and os.path.exists(chart_img):
                self._replace_with_picture(rev_slide, "Chart 4", chart_img)

        # 9. STORE_STOCK_INVENTORY — ẩn với báo cáo khai trương (đồng lý do với STORE_REVENUE).
        stock_slide = slide_map.get("STORE_STOCK_INVENTORY")
        if stock_slide and is_opening_report:
            slides_to_delete.append(stock_slide)
        elif stock_slide:
            if missing_stock:
                self._fill_text(stock_slide, "KPI_STOCK_TOTAL", "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_NGUYEN_GIA", "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_SALE", "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_THANH_LY", "-", font_size=18, bold=True, align="center")
            else:
                self._fill_text(stock_slide, "KPI_STOCK_TOTAL", f"{data.stock.total_qty:,.0f} sản phẩm\n({data.stock.skus_count:,.0f} SKU)" if data.stock.total_qty is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_NGUYEN_GIA", f"{data.stock.qty_nguyen_gia:,.0f} sản phẩm" if data.stock.qty_nguyen_gia is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_SALE", f"{data.stock.qty_sale:,.0f} sản phẩm" if data.stock.qty_sale is not None else "-", font_size=18, bold=True, align="center")
                self._fill_text(stock_slide, "KPI_STOCK_THANH_LY", f"{data.stock.qty_thanh_ly:,.0f} sản phẩm" if data.stock.qty_thanh_ly is not None else "-", font_size=18, bold=True, align="center")
                
                # Set Action Title
                qty_1y = data.stock.age_groups.get("Hàng nguyên giá PP > 1 năm", 0)
                total_valid = sum(data.stock.age_groups.values()) or 1
                pct_1y = (qty_1y / total_valid) * 100
                self._set_slide_title(stock_slide, "3.2.", f"3.2. CƠ CẤU TỒN KHO - HÀNG NGUYÊN GIÁ TRÊN 1 NĂM CHIẾM {pct_1y:.1f}%")

                # Re-position TextBox 1 to make space for stock chart next to it
                for s in stock_slide.shapes:
                    if s.name == "TextBox 1":
                        s.left = int(0.83 * 914400)
                        s.top = int(4.5 * 914400)
                        s.width = int(10.0 * 914400)
                        s.height = int(5.5 * 914400)
                
                # Overwrite age inventory placeholder (TextBox 1)
                self._fill_text(stock_slide, "TextBox 1", data.stock.commentary or "-", font_size=14, bold=False, align="left")
                
                # Add doughnut chart picture next to the commentary box
                stock_chart_img = temp_image_paths.get("stock_chart", "")
                if stock_chart_img and os.path.exists(stock_chart_img):
                    stock_slide.shapes.add_picture(
                        stock_chart_img, 
                        int(11.5 * 914400), 
                        int(4.5 * 914400), 
                        int(7.5 * 914400), 
                        int(5.5 * 914400)
                    )

        # 10. STORE_BEST_SELLERS — ẩn với báo cáo khai trương (chưa có lịch sử bán hàng ý nghĩa).
        best_slide = slide_map.get("STORE_BEST_SELLERS")
        if best_slide and is_opening_report:
            slides_to_delete.append(best_slide)
        elif best_slide:
            self._draw_product_grid(best_slide, "TBL_BEST_SELLERS", data.best_sellers, is_best_seller=True)

        # 11. STORE_SLOW_SELLERS — ẩn với báo cáo khai trương (đồng lý do).
        slow_slide = slide_map.get("STORE_SLOW_SELLERS")
        if slow_slide and is_opening_report:
            slides_to_delete.append(slow_slide)
        elif slow_slide:
            self._draw_product_grid(slow_slide, "TBL_SLOW_SELLERS", data.slow_sellers, is_best_seller=False)

        # 12. STORE_STAFF_LIST (Slide 14) - In-place population
        staff_slide = slide_map.get("STORE_STAFF_LIST")
        if staff_slide:
            table_shape = None
            for s in staff_slide.shapes:
                if s.name == "TBL_STORE_STAFF" and s.has_table:
                    table_shape = s
                    break
            if table_shape:
                table = table_shape.table
                staff_list = list(data.staff.staff_list)
                
                # Sort: CHT first, CHP second, then others alphabetically by name
                def staff_sort_key(item):
                    r = item.role.lower()
                    if "trưởng" in r or "cht" in r:
                        return (0, item.name)
                    elif "phó" in r or "chp" in r:
                        return (1, item.name)
                    else:
                        return (2, item.name)
                staff_list.sort(key=staff_sort_key)
                
                num_staff = len(staff_list)
                if num_staff == 0:
                    self._set_cell_text_preserve_format(table.cell(1, 0), "Không có dữ liệu nhân sự cửa hàng trong hệ thống.")
                    for col in range(1, 5):
                        self._set_cell_text_preserve_format(table.cell(1, col), "")
                    for row in range(2, 7):
                        for col in range(5):
                            self._set_cell_text_preserve_format(table.cell(row, col), "")
                else:
                    for idx in range(6):
                        row_num = idx + 1
                        if idx < num_staff:
                            item = staff_list[idx]
                            notes = item.notes or ""
                            if row_num == 6 and num_staff > 6:
                                notes = f"{notes}; Còn {num_staff - 6} nhân sự khác."
                            
                            self._set_cell_text_preserve_format(table.cell(row_num, 0), item.name)
                            self._set_cell_text_preserve_format(table.cell(row_num, 1), item.role)
                            self._set_cell_text_preserve_format(table.cell(row_num, 2), f"{item.seniority:.1f} năm" if item.seniority is not None else "-")
                            self._set_cell_text_preserve_format(table.cell(row_num, 3), item.skill_rating or "Đạt")
                            self._set_cell_text_preserve_format(table.cell(row_num, 4), notes)
                        else:
                            for col in range(5):
                                self._set_cell_text_preserve_format(table.cell(row_num, col), "")

        # 12.5. STORE_SURVEY_SUPPORT (Slide 16) - In-place population
        survey_slide = slide_map.get("STORE_SURVEY_SUPPORT")
        if survey_slide:
            survey_dict = c_data.get("survey", {})
            for shape in survey_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    lines = [line.strip() for line in shape.text_frame.text.split("\n") if line.strip()]
                    if lines and re.match(r"^[A-E]\d\.", lines[0]):
                        questions = []
                        for line in lines:
                            match = re.match(r"^([A-E]\d)\.\s*(.*)", line)
                            if match:
                                q_key = match.group(1)
                                q_text = match.group(2)
                                questions.append((q_key, q_text))
                        if questions:
                            shape.text_frame.text = ""
                            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                            for i, (q_key, q_text) in enumerate(questions):
                                p = shape.text_frame.add_paragraph() if i > 0 else shape.text_frame.paragraphs[0]
                                p.space_after = Pt(4)
                                
                                # Add question run
                                run_q = p.add_run()
                                run_q.text = f"{q_key}. {q_text}: "
                                run_q.font.name = FONT_PRIMARY
                                run_q.font.size = Pt(8.5)
                                run_q.font.bold = True
                                run_q.font.color.rgb = RGBColor(10, 35, 66)
                                
                                # Add answer run
                                run_a = p.add_run()
                                raw_ans = survey_dict.get(q_key)
                                if isinstance(raw_ans, dict):
                                    raw_ans = raw_ans.get("answer", raw_ans.get("value", ""))
                                ans_str = str(raw_ans).strip() if raw_ans is not None else ""
                                if not ans_str or ans_str == "[object Object]" or ans_str.lower() in ["none", "nan", "null"]:
                                    ans_str = "Chưa trả lời"
                                run_a.text = ans_str
                                run_a.font.name = FONT_PRIMARY
                                run_a.font.size = Pt(8.5)
                                run_a.font.italic = True
                                run_a.font.color.rgb = RGBColor(0, 128, 128)

        # 13. STORE_PENDING_ISSUES (Slide 17) - In-place population
        issues_slide = slide_map.get("STORE_PENDING_ISSUES")
        if issues_slide:
            table_shape = None
            for s in issues_slide.shapes:
                if s.name == "TBL_PENDING_ISSUES" and s.has_table:
                    table_shape = s
                    break
            if table_shape:
                table = table_shape.table
                issues_list = list(data.issues)
                num_issues = len(issues_list)
                if num_issues == 0:
                    self._set_cell_text_preserve_format(table.cell(1, 0), "-")
                    self._set_cell_text_preserve_format(table.cell(1, 1), "Không ghi nhận vấn đề tồn đọng tại thời điểm kiểm tra.")
                    self._set_cell_text_preserve_format(table.cell(1, 2), "-")
                    self._set_cell_text_preserve_format(table.cell(1, 3), "-")
                    self._set_cell_text_preserve_format(table.cell(1, 4), "-")
                    self._set_cell_text_preserve_format(table.cell(1, 5), "-")
                    self._set_cell_text_preserve_format(table.cell(1, 6), "-")
                    
                    for r in (2, 3):
                        for col in range(7):
                            self._set_cell_text_preserve_format(table.cell(r, col), "")
                else:
                    for idx in range(3):
                        row_num = idx + 1
                        if idx < num_issues:
                            item = issues_list[idx]
                            
                            # Format date to DD/MM/YY
                            item_date_short = "-"
                            if item.date:
                                match_dt = re.match(r"^(\d{1,2})[-/](\d{1,2})[-/](\d{4})$", item.date.strip())
                                if match_dt:
                                    d, m, y = match_dt.groups()
                                    item_date_short = f"{int(d):02d}/{int(m):02d}/{y[2:]}"
                                else:
                                    match_dt2 = re.match(r"^(\d{1,2})[-/](\d{1,2})[-/](\d{2})$", item.date.strip())
                                    if match_dt2:
                                        d, m, y = match_dt2.groups()
                                        item_date_short = f"{int(d):02d}/{int(m):02d}/{y}"
                                    else:
                                        item_date_short = item.date.strip()

                            notes = item.notes or ""
                            deadline = "-"
                            clean_notes_parts = []
                            if notes:
                                parts = [p.strip() for p in notes.split("|")]
                                for p in parts:
                                    if p.startswith("Hạn:"):
                                        date_val = p.replace("Hạn:", "").strip()
                                        if date_val:
                                            match_d = re.match(r"^(\d{1,2})[-/](\d{1,2})[-/](\d{4})$", date_val)
                                            if match_d:
                                                d, m, y = match_d.groups()
                                                deadline = f"{int(d):02d}/{int(m):02d}/{y[2:]}"
                                            else:
                                                match_d2 = re.match(r"^(\d{1,2})[-/](\d{1,2})[-/](\d{2})$", date_val)
                                                if match_d2:
                                                    d, m, y = match_d2.groups()
                                                    deadline = f"{int(d):02d}/{int(m):02d}/{y}"
                                                else:
                                                    deadline = date_val
                                    elif p.startswith("Kế hoạch:"):
                                        continue
                                    elif p.startswith("Độ ưu tiên:"):
                                        clean_notes_parts.append(p)
                                    else:
                                        if "kế hoạch" not in p.lower() and "hạn" not in p.lower():
                                            clean_notes_parts.append(p)
                            
                            clean_notes_str = " | ".join(clean_notes_parts)
                            if row_num == 3 and num_issues > 3:
                                separator = " | " if clean_notes_str else ""
                                clean_notes_str = f"{clean_notes_str}{separator}Còn {num_issues - 3} vấn đề khác."
                                
                            if not clean_notes_str:
                                clean_notes_str = "-"
                                
                            self._set_cell_text_preserve_format(table.cell(row_num, 0), str(item.index))
                            self._set_cell_text_preserve_format(table.cell(row_num, 1), item.issue)
                            self._set_cell_text_preserve_format(table.cell(row_num, 2), item_date_short)
                            self._set_cell_text_preserve_format(table.cell(row_num, 3), item.assignee or "-")
                            self._set_cell_text_preserve_format(table.cell(row_num, 4), deadline)
                            status_str = item.status or "Chưa xử lý"
                            cell_status = table.cell(row_num, 5)
                            self._set_cell_text_preserve_format(cell_status, status_str)
                            try:
                                cell_status.fill.solid()
                                if "chưa" in status_str.lower():
                                    cell_status.fill.fore_color.rgb = RGBColor(240, 206, 206) # soft red
                                elif "đang" in status_str.lower():
                                    cell_status.fill.fore_color.rgb = RGBColor(255, 242, 204) # soft yellow
                                elif "đã" in status_str.lower():
                                    cell_status.fill.fore_color.rgb = RGBColor(226, 240, 217) # soft green
                            except Exception as fill_err:
                                print(f"Error coloring status cell: {fill_err}")
                            self._set_cell_text_preserve_format(table.cell(row_num, 6), clean_notes_str)
                        else:
                            for col in range(7):
                                self._set_cell_text_preserve_format(table.cell(row_num, col), "")

            # Operational & CSVC Merge logic
            if merge_operational_csvc:
                # 1. Remove CSVC photo/text shapes at the bottom-left
                self._remove_shape_by_name(issues_slide, "PIC_CSVC_ISSUE")
                self._remove_shape_by_name(issues_slide, "TXT_CSVC_IMAGE_PLACEHOLDER")
                self._remove_shape_by_name(issues_slide, "TXT_CSVC_ISSUE_NOTE")
                
                # Remove CSVC warning label
                for s in list(issues_slide.shapes):
                    if s.has_text_frame and "Vấn đề CSVC cần báo" in s.text_frame.text:
                        try: issues_slide.shapes._spTree.remove(s._element)
                        except Exception: pass
                        
                # 2. Add 3 neat rounded rectangular cards and text columns side by side
                card_w = Inches(2.8)
                card_h = Inches(2.8)
                top_y = Inches(6.8)
                
                lefts = [Inches(0.75), Inches(3.75), Inches(6.75)]
                titles = ["1. LỖI KHO & PHÒNG THỬ", "2. LỖI QUẦY THU NGÂN & BẢO VỆ", "3. VẤN ĐỀ CƠ SỞ VẬT CHẤT"]
                contents = [warehouse_comment, cashier_comment, csvc_comment]
                
                for idx in range(3):
                    left_x = lefts[idx]
                    
                    # Add grey background card
                    card = issues_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, card_w, card_h)
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(245, 247, 250)
                    card.line.color.rgb = RGBColor(220, 225, 230)
                    card.line.width = Pt(1)
                    
                    # Add textbox inside card
                    tb = issues_slide.shapes.add_textbox(left_x + Inches(0.1), top_y + Inches(0.1), card_w - Inches(0.2), card_h - Inches(0.2))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    
                    p_title = tf.paragraphs[0]
                    p_title.text = titles[idx]
                    p_title.font.name = FONT_PRIMARY
                    p_title.font.size = Pt(10)
                    p_title.font.bold = True
                    p_title.font.color.rgb = RGBColor(12, 35, 64)
                    p_title.space_after = Pt(6)
                    
                    p_content = tf.add_paragraph()
                    p_content.text = contents[idx]
                    p_content.font.name = FONT_PRIMARY
                    p_content.font.size = Pt(8.5)
                    p_content.font.color.rgb = RGBColor(0, 0, 0)
            else:
                # Normal/Adaptive CSVC photo mapping
                self._fill_text(issues_slide, "TXT_CSVC_ISSUE_NOTE", csvc_comment)
                
                if csvc_path and os.path.exists(csvc_path) and os.path.getsize(csvc_path) > 0:
                    self._fill_image_slot(issues_slide, "PIC_CSVC_ISSUE", "TXT_CSVC_IMAGE_PLACEHOLDER", csvc_path)
                else:
                    self._remove_shape_by_name(issues_slide, "PIC_CSVC_ISSUE")
                    self._remove_shape_by_name(issues_slide, "TXT_CSVC_IMAGE_PLACEHOLDER")
                    for s in issues_slide.shapes:
                        if s.name == "TXT_CSVC_ISSUE_NOTE":
                            s.height = Inches(2.44)
                            if s.has_text_frame:
                                s.text_frame.vertical_anchor = MSO_ANCHOR.TOP

            # Action Plan proposal
            action_plan_text = "Chưa ghi nhận đề xuất khắc phục bổ sung tại thời điểm kiểm tra."
            if data.form_response and data.form_response.action_plan:
                action_plan_text = data.form_response.action_plan
            elif data.frontage_action:
                action_plan_text = data.frontage_action
            self._fill_text(issues_slide, "TXT_ACTION_PLAN_PROPOSALS", action_plan_text)

        # 14. STORE_DEV_PROPOSALS (Slide 17)
        dev_slide = slide_map.get("STORE_DEV_PROPOSALS")
        if dev_slide:
            # Thu nhập dữ liệu cửa hàng để làm ngữ cảnh cho AI/Fallback
            store_name = data.metadata.store_name if data.metadata else "-"
            asm_name = data.form_response.asm_name if data.form_response and data.form_response.asm_name else (data.metadata.asm_name if data.metadata else "QLKD")
            
            # Doanh thu
            rev_actual = data.revenue.revenue_actual if data.revenue else 0
            rev_target = data.revenue.revenue_target if data.revenue else 0
            attainment = data.revenue.attainment_pct if data.revenue else 0.0
            rev_comment = data.revenue.commentary if data.revenue else ""
            
            # Tồn kho
            stock_total = data.stock.total_qty if data.stock else 0
            stock_ng = data.stock.qty_nguyen_gia if data.stock else 0
            stock_sale = data.stock.qty_sale if data.stock else 0
            stock_tl = data.stock.qty_thanh_ly if data.stock else 0
            stock_comment = data.stock.commentary if data.stock else ""
            
            # Lỗi vận hành ghi nhận
            issues_list = []
            if data.issues:
                for issue in data.issues:
                    issues_list.append(f"- {issue.label}: {issue.issue} (Trạng thái: {issue.status})")
            issues_text = "\n".join(issues_list) if issues_list else "Không ghi nhận lỗi vận hành nghiêm trọng."
            
            # Nhận xét từ người kiểm tra
            front_comment = ""
            inner_comment = ""
            human_rec = ""
            if data.form_response:
                front_comment = data.form_response.comment_frontage or ""
                inner_comment = data.form_response.comment_inner or data.form_response.comment_merch or ""
                human_rec = data.form_response.store_recommendation or ""
            
            # Lấy sản phẩm bán chạy động từ dữ liệu để tránh hardcode
            best_sellers_list = []
            if data.best_sellers:
                for item in data.best_sellers[:2]:
                    if item.product_name:
                        name = item.product_name.strip()
                        if name:
                            best_sellers_list.append(name)
            
            if best_sellers_list:
                if len(best_sellers_list) == 1:
                    best_seller_phrase = f"mặt hàng {best_sellers_list[0]}"
                else:
                    best_seller_phrase = f"các mặt hàng như {best_sellers_list[0]}, {best_sellers_list[1]}"
            else:
                best_seller_phrase = "các nhóm hàng bán chạy thực tế tại chi nhánh"

            # Lấy sản phẩm bán chậm động từ dữ liệu
            slow_sellers_list = []
            if data.slow_sellers:
                for item in data.slow_sellers[:2]:
                    if item.product_name:
                        name = item.product_name.strip()
                        if name:
                            slow_sellers_list.append(name)
            
            if slow_sellers_list:
                if len(slow_sellers_list) == 1:
                    slow_seller_phrase = f"mặt hàng {slow_sellers_list[0]}"
                else:
                    slow_seller_phrase = f"các mặt hàng như {slow_sellers_list[0]}, {slow_sellers_list[1]}"
            else:
                slow_seller_phrase = "các nhóm hàng bán chậm tại chi nhánh"

            # Xây dựng đề xuất offline (rule-based fallback) động chuyên sâu
            fallback_rev_bullets = []
            fallback_merch_bullets = []
            
            # Phần doanh thu (tối đa 2 bullet, mỗi bullet dưới 22 từ)
            if attainment < 85:
                fallback_rev_bullets.append(f"- Đẩy mạnh trưng bày, tư vấn bán {best_seller_phrase} để cải thiện nhanh doanh thu.")
                fallback_rev_bullets.append("- Chủ động kết nối khách hàng VIP chi nhánh qua Zalo/SMS giới thiệu sản phẩm mới về.")
            else:
                fallback_rev_bullets.append("- Duy trì tăng trưởng tốt, bán chéo phụ kiện thời trang để nâng cao chỉ số UPT.")
                fallback_rev_bullets.append("- Tối ưu trưng bày mặt tiền thu hút thêm lượt khách hàng vãng lai vào chi nhánh.")
                
            # Phần hàng hóa (tối đa 2 bullet, mỗi bullet dưới 22 từ)
            if stock_tl > 50 or stock_sale > 100:
                fallback_merch_bullets.append("- Luân chuyển các mã hàng lẻ size, chậm bán hoặc thanh lý sang điểm bán tốt hơn.")
                fallback_merch_bullets.append("- Trưng bày hàng giảm giá đúng phân khu riêng biệt để giải phóng nhanh tồn kho lâu ngày.")
            elif data.issues:
                first_issue_label = data.issues[0].label if data.issues else "trưng bày"
                fallback_merch_bullets.append(f"- Tập trung khắc phục triệt để lỗi {first_issue_label.lower()} và vệ sinh kệ đã ghi nhận tại quầy.")
                fallback_merch_bullets.append("- Theo dõi tiến độ xử lý các lỗi vận hành phát hiện và báo cáo khắc phục cho ASM.")
            else:
                fallback_merch_bullets.append("- Sắp xếp kho bãi ngăn nắp, kiểm kho định kỳ để tránh lệch size số và thất thoát hàng.")
                fallback_merch_bullets.append("- Đối chiếu cơ cấu hàng nguyên giá và hàng sale để tối ưu diện tích trưng bày quầy kệ.")
                
            rule_based_fallback = (
                "ĐỀ XUẤT PHÁT TRIỂN DOANH THU:\n"
                f"{fallback_rev_bullets[0]}\n"
                f"{fallback_rev_bullets[1]}\n"
                "GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA:\n"
                f"{fallback_merch_bullets[0]}\n"
                f"{fallback_merch_bullets[1]}"
            )

            # Khởi động AIClient
            dev_rec = None
            ai_client = None
            try:
                # Fix 30-07: "modules.ai_client" không tồn tại trong repo này (thư mục modules/
                # không có) — ImportError bị nuốt lặng lẽ ở except bên dưới, khiến slide này
                # LUÔN dùng fallback rule-based dù cascade AI thật đã có sẵn ở reports/ai_client.py
                # (cùng chữ ký generate(prompt, fallback=...) + .last_source, drop-in tương thích).
                from reports.ai_client import AIClient
                ai_client = AIClient()
            except Exception as e:
                print(f"[AI Recommendation] Warning: Cannot initialize AIClient: {e}")

            if ai_client:
                prompt = f"""
Bạn là một Giám đốc Kinh doanh Vùng (ASM) chuyên nghiệp tại chuỗi thời trang An Phước.
Dựa vào dữ liệu thực tế đợt kiểm tra cửa hàng dưới đây, hãy tự động đưa ra các đề xuất hành động cụ thể, thiết thực và mang tính chiến đấu cao.

THÔNG TIN CỬA HÀNG:
- Cửa hàng: An Phước {store_name}
- ASM: {asm_name}

DỮ LIỆU DOANH THU (MTD):
- Doanh thu: {rev_actual:,.0f} VNĐ / Mục tiêu: {rev_target:,.0f} VNĐ (Đạt: {attainment:.1f}%)
- Nhận xét doanh thu: {rev_comment}

DỮ LIỆU TỒN KHO:
- Tổng tồn: {stock_total:,.0f} SP (Nguyên giá: {stock_ng:,.0f} | Sale: {stock_sale:,.0f} | Thanh lý: {stock_tl:,.0f})
- Phân tích tồn kho theo đợt phân phối: {stock_comment}

GHI NHẬN THỰC TẾ:
- Sản phẩm bán chạy nhất: {best_seller_phrase}
- Sản phẩm bán chậm nhất: {slow_seller_phrase}
- Lỗi vận hành ghi nhận:
{issues_text}
- Nhận xét mặt tiền ngoại quan: {front_comment}
- Nhận xét trưng bày bên trong: {inner_comment}
- Ý kiến ghi nhận của ASM: {human_rec}

YÊU CẦU ĐỊNH DẠNG CỰC KỲ NGHIÊM NGẶT:
Đoạn văn viết bằng tiếng Việt, chia làm đúng 2 phần với tiêu đề in hoa rõ ràng:
ĐỀ XUẤT PHÁT TRIỂN DOANH THU:
- [Bullet 1 hành động cụ thể, không quá 22 từ]
- [Bullet 2 hành động cụ thể, không quá 22 từ]
GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA:
- [Bullet 3 hành động cụ thể, không quá 22 từ]
- [Bullet 4 hành động cụ thể, không quá 22 từ]

Quy định viết để đảm bảo chất lượng kiểm soát (QC) và tuyệt đối không Hallucination:
- BẮT BUỘC có đúng 2 tiêu đề trên.
- Tổng số bullet trong toàn bộ văn bản phải BẮT BUỘC từ 2 đến 4 bullet (mỗi phần 1-2 bullet).
- Mỗi bullet BẮT BUỘC dài tối đa 22 từ, viết ngắn gọn, súc tích, đi thẳng vào hành động thực tế.
- Bám sát số liệu doanh thu và tồn kho ở trên.
- TUYỆT ĐỐI KHÔNG BỊA ĐẶT các chương trình khuyến mãi (như 'tặng voucher', 'giảm giá 50%', 'mua 1 tặng 1'...) nếu đề bài không đề cập.
- TUYỆT ĐỐI KHÔNG BỊA ĐẶT các kênh tiếp thị ngoài cửa hàng như 'chạy quảng cáo facebook', 'phát tờ rơi', 'treo băng rôn'.
- Chỉ đưa ra các hành động thực tế tại cửa hàng: ví dụ nếu hàng cũ phân phối trên 1 năm chiếm tỷ trọng cao (>50%), đề xuất tập trung trưng bày/luân chuyển nhóm này. Nếu doanh thu thấp, tập trung thúc đẩy sản phẩm bán chạy hoặc đào tạo nhân sự.
- Chỉ trả về nội dung đề xuất trực tiếp, không viết lời dẫn hay giải thích gì thêm.
"""
                try:
                    print(f"[AI Recommendation] Requesting AI recommendation for store {store_name}...")
                    ai_response = ai_client.generate(prompt, fallback="")
                    
                    # Clean and trim the AI response first
                    trimmed_response = self._clean_and_trim_proposals(ai_response)
                    
                    # Validate the trimmed response with strict QC
                    if self._validate_proposals_no_hallucination(trimmed_response):
                        print(f"[AI Recommendation] AI response validated and trimmed successfully (source: {ai_client.last_source})")
                        dev_rec = trimmed_response
                    else:
                        print("[AI Recommendation] AI response failed strict QC validation after trimming, falling back to rule-based proposals...")
                except Exception as ai_err:
                    print(f"[AI Recommendation] Error calling AI: {ai_err}")
            
            if not dev_rec:
                print("[AI Recommendation] Using dynamic rule-based fallback proposals.")
                dev_rec = rule_based_fallback

            # Ensure the final output is trimmed and formatted (handles fallback trimming as well)
            dev_rec = self._clean_and_trim_proposals(dev_rec)

            self._fill_text(dev_slide, "TXT_DEV_PROPOSALS", dev_rec)

        # 15. STORE_THANK_YOU (Slide 18)
        thank_you_slide = slide_map.get("STORE_THANK_YOU")
        if thank_you_slide:
            report_date = data.form_response.report_date if data.form_response and data.form_response.report_date else "-"
            asm_name = data.form_response.asm_name if data.form_response and data.form_response.asm_name else data.metadata.asm_name
            footer_text = f"Cửa hàng: An Phước {data.metadata.store_name}   |   Ngày: {report_date}   |   QLKD: {asm_name}"
            self._fill_text(thank_you_slide, "TXT_FOOTER_INFO", footer_text, font_size=14, align="center")

        # Delete collected slides
        for slide in slides_to_delete:
            self._delete_slide(prs, slide)
            
        # Clean up unused placeholder shapes and overlay texts from remaining slides
        self._cleanup_unused_placeholders(prs)

        # Chuẩn hoá font toàn bộ về Be Vietnam Pro (diệt font lạ như Abadi/Inter còn sót
        # ở text tĩnh của template) → hết tình trạng "font không đều".
        self._normalize_fonts(prs)

        # Save presentation
        prs.save(output_path)
        print(f"Generated individual store report: {output_path}")

    # Font ký hiệu/icon — KHÔNG đổi (giữ nguyên để không vỡ icon).
    _SYMBOL_FONTS = ("wingding", "webding", "symbol", "segoe mdl", "segoe fluent",
                     "material", "font awesome", "glyphicon")

    def _normalize_fonts(self, prs):
        """Đặt mọi run text về FONT_PRIMARY, trừ font ký hiệu. Bao gồm cả bảng & group shapes."""
        def _fix_tf(tf):
            for p in tf.paragraphs:
                for r in p.runs:
                    fn = (r.font.name or "").lower()
                    if fn and any(sym in fn for sym in self._SYMBOL_FONTS):
                        continue
                    r.font.name = FONT_PRIMARY

        def _walk(shapes):
            for sh in shapes:
                try:
                    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _walk(sh.shapes)
                        continue
                except Exception:
                    pass
                if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                    _fix_tf(sh.text_frame)
                if getattr(sh, "has_table", False) and sh.has_table:
                    for row in sh.table.rows:
                        for cell in row.cells:
                            _fix_tf(cell.text_frame)

        for slide in prs.slides:
            _walk(slide.shapes)

    def _get_slide_id(self, slide) -> str:
        for shape in slide.shapes:
            if shape.name == "META_SLIDE_ID" and shape.has_text_frame:
                return shape.text_frame.text.strip()
        return ""

    def _clean_and_trim_proposals(self, text: str) -> str:
        """
        Clean, trim, and format AI-generated recommendations:
        - Must have two sections: 'ĐỀ XUẤT PHÁT TRIỂN DOANH THU:' and 'GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA:'.
        - Limits to at most 2 bullets per section (total max 4 bullets).
        - Truncates each bullet to at most 22 words and ensures it ends with a period.
        """
        if not text:
            return ""
        
        lines = text.split("\n")
        section1_bullets = []
        section2_bullets = []
        current_section = 0  # 1 for section 1, 2 for section 2
        
        for line in lines:
            line_strip = line.strip()
            if not line_strip:
                continue
            
            # Detect section headings
            line_upper = line_strip.upper().replace(":", "")
            if "ĐỀ XUẤT PHÁT TRIỂN DOANH THU" in line_upper:
                current_section = 1
                continue
            elif "GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA" in line_upper:
                current_section = 2
                continue
            
            if current_section in (1, 2):
                # Clean up bullet markers and formatting
                bullet_content = line_strip.lstrip("-*•◦ ").strip()
                bullet_content = bullet_content.replace("**", "")
                if not bullet_content:
                    continue
                
                # Trim to max 22 words
                words = bullet_content.split()
                if len(words) > 22:
                    words = words[:22]
                    bullet_content = " ".join(words)
                    if not bullet_content.endswith((".", "!", "?")):
                        bullet_content += "."
                
                if current_section == 1:
                    section1_bullets.append(bullet_content)
                else:
                    section2_bullets.append(bullet_content)
        
        # Keep at most 2 bullets per section
        section1_bullets = section1_bullets[:2]
        section2_bullets = section2_bullets[:2]
        
        # Reconstruct structured text
        result_parts = []
        if section1_bullets:
            result_parts.append("ĐỀ XUẤT PHÁT TRIỂN DOANH THU:")
            for b in section1_bullets:
                result_parts.append(f"- {b}")
        if section2_bullets:
            result_parts.append("GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA:")
            for b in section2_bullets:
                result_parts.append(f"- {b}")
                
        return "\n".join(result_parts)

    def _validate_proposals(self, text: str) -> bool:
        """Verify that the proposals meet all length, heading, and bullet count requirements."""
        if not text:
            print("[AI Validation] Failed: Empty text")
            return False
        
        text_upper = text.upper()
        if "ĐỀ XUẤT PHÁT TRIỂN DOANH THU" not in text_upper:
            print("[AI Validation] Failed: Missing heading 'ĐỀ XUẤT PHÁT TRIỂN DOANH THU'")
            return False
        if "GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA" not in text_upper:
            print("[AI Validation] Failed: Missing heading 'GIẢI QUYẾT VẤN ĐỀ HÀNG HÓA'")
            return False
        
        lines = text.split("\n")
        bullets = [l.strip() for l in lines if l.strip().startswith("-")]
        
        if len(bullets) < 2 or len(bullets) > 4:
            print(f"[AI Validation] Failed: bullet count is {len(bullets)} (must be 2-4)")
            return False
            
        for i, b in enumerate(bullets):
            content = b.lstrip("- ").strip()
            words = content.split()
            word_count = len(words)
            if word_count > 22:
                print(f"[AI Validation] Failed: bullet {i+1} has {word_count} words (max 22): '{content}'")
                return False
            if word_count < 3:
                print(f"[AI Validation] Failed: bullet {i+1} is too short ({word_count} words): '{content}'")
                return False
                
        return True


    def _validate_proposals_no_hallucination(self, text: str) -> bool:
        """Strict QC verification to prevent AI hallucination of made-up promos or ads."""
        if not self._validate_proposals(text):
            return False
            
        forbidden_keywords = [
            "voucher", "tặng quà", "mua 1 tặng 1", "giảm giá 50%", "tri ân khách hàng",
            "quảng cáo facebook", "chạy quảng cáo", "tờ rơi", "băng rôn", "tài trợ",
            "social media", "marketing online", "quay clip", "livestream", "tiktok"
        ]
        
        text_lower = text.lower()
        for kw in forbidden_keywords:
            if kw in text_lower:
                print(f"[AI Validation QC] Failed: Hallucination keyword '{kw}' found in response.")
                return False
                
        return True


    def _fill_text(self, slide, shape_name: str, new_text: str, font_size: int = 14, bold: bool = False, align: str = "left"):
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                # 1. Detect original style from template
                font_name = FONT_PRIMARY
                template_font_size = None
                text_color = CLR_INK

                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    # Font: luôn chuẩn hoá về FONT_PRIMARY để hết tình trạng font không đều.
                    if run.font.size:
                        template_font_size = run.font.size.pt
                    try:
                        # Chỉ kế thừa màu template nếu KHÔNG phải xám placeholder mờ.
                        if run.font.color and run.font.color.rgb and not _is_placeholder_gray(run.font.color.rgb):
                            text_color = run.font.color.rgb
                    except (AttributeError, TypeError):
                        pass
                
                # Use template font size if caller passed the default 14
                if font_size == 14 and template_font_size is not None:
                    base_size = template_font_size
                else:
                    base_size = font_size
                
                # 2. Dynamic font size adjustment based on text length to prevent overflow
                text_len = len(str(new_text))
                actual_font_size = base_size
                if base_size <= 18: # scale down slightly for body/commentary text if extremely long
                    if text_len > 400:
                        actual_font_size = min(actual_font_size, 12.0)
                    elif text_len > 250:
                        actual_font_size = min(actual_font_size, 13.0)
                    elif text_len > 120:
                        actual_font_size = min(actual_font_size, 14.0)
                
                # 3. Clear text frame and enable word wrap + narrow margins
                tf = shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.05)
                tf.margin_right = Inches(0.05)
                tf.margin_top = Inches(0.05)
                tf.margin_bottom = Inches(0.05)
                tf.clear()
                
                # 4. Fill text and style each paragraph individually
                lines = str(new_text).split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = line
                    
                    # Set alignment
                    if align == "center":
                        p.alignment = PP_ALIGN.CENTER
                    elif align == "right":
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT
                    
                    p.font.name = font_name
                    p.font.size = Pt(actual_font_size)
                    p.font.bold = bold
                    p.font.color.rgb = text_color
                return
        print(f"Warning: Text shape '{shape_name}' not found on slide.")

    def _fill_opening_info(self, slide, form_response):
        """Đồng bộ 30-07: thay TXT_GENERAL_COMMENT bằng thông tin khai trương (loại/giai
        đoạn/ngày) + badge mức độ sẵn sàng màu theo CLR_OK/CLR_WARN/CLR_ERR, dùng cho báo
        cáo inspection_mode == "opening". Ghi trực tiếp run/paragraph (không qua _fill_text)
        để tô màu riêng cho dòng mức độ sẵn sàng."""
        OPENING_TYPE_LABELS = {"new": "Mở mới", "reopen": "Tái khai trương"}
        OPENING_PHASE_LABELS = {"before": "Trước khai trương", "day": "Ngày khai trương", "after": "Sau khai trương"}
        READINESS_LABELS = {"ready": "SẴN SÀNG", "minor_fix": "CẦN KHẮC PHỤC NHỎ", "not_ready": "CHƯA SẴN SÀNG"}
        READINESS_COLORS = {"ready": CLR_OK, "minor_fix": CLR_WARN, "not_ready": CLR_ERR}

        opening_type = getattr(form_response, "opening_type", None) if form_response else None
        opening_phase = getattr(form_response, "opening_phase", None) if form_response else None
        opening_date = getattr(form_response, "opening_date", None) if form_response else None
        opening_readiness = getattr(form_response, "opening_readiness", None) if form_response else None

        type_label = OPENING_TYPE_LABELS.get(opening_type, opening_type or "-")
        phase_label = OPENING_PHASE_LABELS.get(opening_phase, opening_phase or "-")
        readiness_label = READINESS_LABELS.get(opening_readiness, opening_readiness or "Chưa xác định")
        readiness_color = READINESS_COLORS.get(opening_readiness, CLR_MUTED)

        for shape in slide.shapes:
            if shape.name == "TXT_GENERAL_COMMENT" and shape.has_text_frame:
                tf = shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.05)
                tf.margin_right = Inches(0.05)
                tf.margin_top = Inches(0.05)
                tf.margin_bottom = Inches(0.05)
                tf.clear()

                p0 = tf.paragraphs[0]
                p0.text = f"🎉 BÁO CÁO KHAI TRƯƠNG — {type_label}"
                p0.font.name = FONT_PRIMARY
                p0.font.size = Pt(14)
                p0.font.bold = True
                p0.font.color.rgb = CLR_NAVY
                p0.space_after = Pt(6)

                p1 = tf.add_paragraph()
                p1.text = f"Giai đoạn: {phase_label}   |   Ngày khai trương chính thức: {opening_date or '-'}"
                p1.font.name = FONT_PRIMARY
                p1.font.size = Pt(11)
                p1.font.color.rgb = CLR_MUTED
                p1.space_after = Pt(10)

                p2 = tf.add_paragraph()
                p2.text = f"MỨC ĐỘ SẴN SÀNG KHAI TRƯƠNG: {readiness_label}"
                p2.font.name = FONT_PRIMARY
                p2.font.size = Pt(14)
                p2.font.bold = True
                p2.font.color.rgb = readiness_color
                return
        print("Warning: Text shape 'TXT_GENERAL_COMMENT' not found on slide (opening info).")

    def _build_opening_before_after(self, slide, form_response):
        """Đồng bộ 30-07: ghi đè 1 slide có sẵn (STORE_INNER_PHOTOS) bằng layout 2 cột
        Trước/Sau sửa chữa cho báo cáo Tái khai trương — tái dùng đúng kỹ thuật "xoá shape
        thân slide rồi vẽ textbox mới" đã chứng minh ở khối merge frontage/inner phía trên."""
        shapes_to_delete = []
        for s in slide.shapes:
            if s.name not in ["META_SLIDE_ID", "META_TEMPLATE_VERSION"] and s.top and s.top > Inches(1.85):
                shapes_to_delete.append(s)
        for s in shapes_to_delete:
            try: slide.shapes._spTree.remove(s._element)
            except Exception: pass

        before_paths = [p.local_path for p in (form_response.photos or [])
                         if p.section == "opening_before" and p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]
        after_paths = [p.local_path for p in (form_response.photos or [])
                        if p.section == "opening_after" and p.local_path and os.path.exists(p.local_path) and os.path.getsize(p.local_path) > 0]

        columns = [("TRƯỚC SỬA CHỮA", before_paths, Inches(1.0)), ("SAU KHI HOÀN THIỆN", after_paths, Inches(10.0))]
        col_w = Inches(8.5)
        for title, paths, left in columns:
            tb = slide.shapes.add_textbox(left, Inches(2.2), col_w, Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_PRIMARY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = CLR_NAVY

            if paths:
                try:
                    slide.shapes.add_picture(paths[0], left, Inches(3.0), col_w, Inches(6.5))
                except Exception as _e:
                    print(f"[opening before/after] Lỗi chèn ảnh {title}: {_e}")
            else:
                ph = slide.shapes.add_textbox(left, Inches(3.0), col_w, Inches(6.5))
                ph_tf = ph.text_frame
                ph_tf.word_wrap = True
                ph_p = ph_tf.paragraphs[0]
                ph_p.text = "(Chưa có ảnh)"
                ph_p.font.name = FONT_PRIMARY
                ph_p.font.size = Pt(12)
                ph_p.font.color.rgb = CLR_MUTED
                ph_p.alignment = PP_ALIGN.CENTER

        # Đổi tiêu đề slide (nếu tìm được shape tiêu đề gốc)
        for s in slide.shapes:
            if s.has_text_frame and ("KHÔNG GIAN" in s.text_frame.text or "BÊN TRONG" in s.text_frame.text):
                s.text_frame.text = "2. SO SÁNH TRƯỚC / SAU SỬA CHỮA"
                if s.text_frame.paragraphs and s.text_frame.paragraphs[0].runs:
                    s.text_frame.paragraphs[0].runs[0].font.name = FONT_PRIMARY
                    s.text_frame.paragraphs[0].runs[0].font.bold = True

    def _fill_text_by_prefix_or_name(self, slide, name_or_prefix: str, new_text: str, font_size: int = 11, bold: bool = False):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if shape.name == name_or_prefix or txt.startswith(name_or_prefix):
                    # 1. Detect original style
                    font_name = FONT_PRIMARY
                    template_font_size = None
                    text_color = RGBColor(0, 0, 0)
                    
                    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                        run = shape.text_frame.paragraphs[0].runs[0]
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            template_font_size = run.font.size.pt
                        try:
                            if run.font.color and run.font.color.rgb:
                                text_color = run.font.color.rgb
                        except AttributeError:
                            pass
                    
                    # Determine base size
                    if font_size == 11 and template_font_size is not None:
                        base_size = template_font_size
                    else:
                        base_size = font_size
                    
                    # 2. Dynamic font size adjustment based on text length
                    text_len = len(str(new_text))
                    actual_font_size = base_size
                    if base_size <= 14:
                        if text_len > 300:
                            actual_font_size = min(actual_font_size, 9.0)
                        elif text_len > 150:
                            actual_font_size = min(actual_font_size, 10.0)
                        elif text_len > 80:
                            actual_font_size = min(actual_font_size, 11.0)
                    
                    # 3. Clear text frame and style
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.05)
                    tf.margin_right = Inches(0.05)
                    tf.margin_top = Inches(0.05)
                    tf.margin_bottom = Inches(0.05)
                    tf.clear()
                    
                    # 4. Populate paragraphs
                    lines = str(new_text).split("\n")
                    for i, line in enumerate(lines):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = line
                        p.font.name = font_name
                        p.font.size = Pt(actual_font_size)
                        p.font.bold = bold
                        p.font.color.rgb = text_color
                    return

    def _build_exec_facts(self, data, c_data) -> list:
        """Thu thập các dữ kiện CÓ THẬT từ báo cáo để làm nguyên liệu cho narrator.
        Chỉ đưa vào dữ liệu tồn tại — không suy diễn, không bịa."""
        facts = []
        # 1. Doanh thu
        try:
            rev = data.revenue
            if rev and (rev.revenue_actual or rev.revenue_target):
                s = f"Doanh thu thực hiện {rev.revenue_actual:,.0f} VNĐ"
                if rev.attainment_pct:
                    s += f", đạt {rev.attainment_pct:.0f}% chỉ tiêu"
                if getattr(rev, 'yoy_change_pct', 0):
                    _d = rev.yoy_change_pct
                    s += f", {'tăng' if _d >= 0 else 'giảm'} {abs(_d):.0f}% so với cùng kỳ năm trước"
                facts.append(s + ".")
        except Exception:
            pass
        # 2. Đánh giá các hạng mục (từ form_response)
        try:
            fr = data.form_response
            if fr:
                rating_map = [
                    ("Mặt tiền", getattr(fr, 'rating_frontage', '')),
                    ("Không gian trong", getattr(fr, 'rating_inner', '')),
                    ("Trưng bày hàng hoá", getattr(fr, 'rating_merch', '')),
                    ("Nhân sự", getattr(fr, 'rating_staff', '')),
                    ("Cơ sở vật chất", getattr(fr, 'rating_csvc', '')),
                ]
                good = [n for n, r in rating_map if r == "Tốt"]
                weak = [n for n, r in rating_map if r == "Chưa đạt"]
                if good:
                    facts.append("Các hạng mục đạt tốt: " + ", ".join(good) + ".")
                if weak:
                    facts.append("Các hạng mục CHƯA đạt cần khắc phục: " + ", ".join(weak) + ".")
        except Exception:
            pass
        # 3. Các lỗi cụ thể (Không đạt) từ checklist_json
        try:
            sections = (c_data or {}).get("sections", {})
            fails = []
            for sec in sections.values():
                for item in (sec.get("items") or []):
                    if item.get("eval") == "Không đạt":
                        lbl = item.get("label", "").strip()
                        note = (item.get("note") or "").strip()
                        if lbl:
                            fails.append(lbl + (f" ({note})" if note else ""))
            for f in fails[:6]:
                facts.append("Lỗi ghi nhận: " + f + ".")
        except Exception:
            pass
        # 4. Tồn kho
        try:
            st = data.stock
            if st and st.total_qty:
                facts.append(f"Tổng tồn kho {st.total_qty:,} sản phẩm "
                             f"(nguyên giá {st.qty_nguyen_gia:,}, sale {st.qty_sale:,}, thanh lý {st.qty_thanh_ly:,}).")
        except Exception:
            pass
        return facts

    def _add_multiline_paragraphs(self, tf, text: str, font_name: str = FONT_PRIMARY, font_size: float = 10.5, color_rgb = RGBColor(0, 0, 0), bold: bool = False):
        """Helper to append multiple paragraphs for multi-line text, styling each paragraph individually."""
        lines = str(text).split("\n")
        # Enable word wrap and narrow margins
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        for i, line in enumerate(lines):
            # If the text frame is empty, we reuse the first paragraph, otherwise add new
            if i == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color_rgb

    def _fill_image_slot(self, slide, pic_name: str, txt_name: str, image_path: str):
        """Helper to fill image containers. Keeps grey fallback intact if missing or invalid."""
        pic_shape = None
        txt_shape = None
        for s in slide.shapes:
            if s.name == pic_name:
                pic_shape = s
            elif s.name == txt_name:
                txt_shape = s
                
        if not pic_shape:
            print(f"Warning: Image container '{pic_name}' not found on slide.")
            return
            
        is_valid_image = False
        if image_path and os.path.exists(image_path) and os.path.getsize(image_path) > 0:
            ext = os.path.splitext(image_path.lower())[1]
            if ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"]:
                is_valid_image = True
                
        if is_valid_image:
            try:
                left = pic_shape.left
                top = pic_shape.top
                width = pic_shape.width
                height = pic_shape.height
                
                # 1. Dynamically fit image to the shape's exact aspect ratio
                temp_dir = os.path.dirname(image_path)
                base = os.path.basename(image_path)
                fitted_path = os.path.join(temp_dir, f"fitted_{width}_{height}_{base}")
                
                # Target width 800, calculate proportional height
                target_w = 800
                target_h = int(800 * height / width) if width > 0 else 600
                
                ImageProcessor.process_and_fit_image(image_path, fitted_path, target_width=target_w, target_height=target_h)
                
                # 2. Insert at same z-order index
                spTree = slide.shapes._spTree
                placeholder_element = pic_shape._element
                insert_index = spTree.index(placeholder_element)
                
                new_picture = slide.shapes.add_picture(fitted_path, left, top, width, height)
                new_pic_element = new_picture._element
                
                # Move element in XML tree
                spTree.remove(new_pic_element)
                spTree.insert(insert_index, new_pic_element)
                
                # Remove placeholder
                spTree.remove(placeholder_element)
                
                # Clear overlay text
                if txt_shape and txt_shape.has_text_frame:
                    txt_shape.text_frame.text = ""
            except Exception as e:
                print(f"Error rendering image {image_path}: {e}")
                is_valid_image = False
                
        if not is_valid_image:
            # Overwrite text shape overlay with fallback text
            if txt_shape and txt_shape.has_text_frame:
                txt_shape.text_frame.text = ""
                p = txt_shape.text_frame.paragraphs[0]
                p.text = "Không có ảnh ghi nhận"
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    run = p.runs[0]
                    run.font.name = FONT_PRIMARY
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                txt_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _fill_image_slots_adaptive(self, slide, slot_configs: list, photo_paths: list):
        """
        Fill image slots adaptively based on how many valid photos are available.

        slot_configs: list of (pic_name, txt_name) tuples matching the photo placeholders in the template.
        photo_paths:  list of local file paths (one per slot, empty string means no photo).

        Behaviour:
          - 0 photos  → hide all PIC_ slots, replace first TXT_ with a centred 'no photo' message.
          - 1 photo   → show in centre slot (or first available), hide others, widen the visible slot.
          - 2 photos  → show in first two slots equally spaced, hide the third.
          - 3 photos  → normal fill (standard layout).
        """
        from pptx.util import Emu

        # Resolve valid (path, slot_config) pairs
        valid_pairs = [(path, cfg) for path, cfg in zip(photo_paths, slot_configs) if path and os.path.exists(path) and os.path.getsize(path) > 0]
        num_valid = len(valid_pairs)

        # Collect all PIC_ shapes and their bounding boxes (for layout calculation)
        all_pic_shapes = {}
        all_txt_shapes = {}
        for (pic_name, txt_name) in slot_configs:
            for s in slide.shapes:
                if s.name == pic_name:
                    all_pic_shapes[pic_name] = s
                elif s.name == txt_name:
                    all_txt_shapes[txt_name] = s

        if num_valid == 0:
            # Hide all image slots; show a single centred message in the first TXT shape
            for (pic_name, txt_name) in slot_configs:
                pic_s = all_pic_shapes.get(pic_name)
                txt_s = all_txt_shapes.get(txt_name)
                if pic_s:
                    try:
                        slide.shapes._spTree.remove(pic_s._element)
                    except Exception:
                        pass
                if txt_s and txt_s.has_text_frame:
                    if txt_name == slot_configs[0][1]:
                        # First slot: show the message spanning all three slots widths
                        try:
                            # Widen to cover all three original slots
                            first_pic = list(all_pic_shapes.values())[0] if all_pic_shapes else None
                            last_pic = list(all_pic_shapes.values())[-1] if len(all_pic_shapes) > 1 else None
                            if first_pic and last_pic:
                                txt_s.left = first_pic.left
                                txt_s.width = (last_pic.left + last_pic.width) - first_pic.left
                        except Exception:
                            pass
                        txt_s.text_frame.text = ""
                        p = txt_s.text_frame.paragraphs[0]
                        p.text = "Không ghi nhận hình ảnh tại khu vực này"
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            run = p.runs[0]
                            run.font.name = FONT_PRIMARY
                            run.font.size = Pt(12)
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        txt_s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    else:
                        # Other TXT slots: blank them
                        try:
                            slide.shapes._spTree.remove(txt_s._element)
                        except Exception:
                            pass
            return

        # For 1 or 2 valid photos: hide unused slots and redistribute space
        if num_valid < len(slot_configs):
            unused_slot_names = set()
            used_pic_names = {cfg[0] for _, cfg in valid_pairs}
            for (pic_name, txt_name) in slot_configs:
                if pic_name not in used_pic_names:
                    unused_slot_names.add(pic_name)
                    unused_slot_names.add(txt_name)

            # Remove unused shapes
            for shape_name in unused_slot_names:
                for s in slide.shapes:
                    if s.name == shape_name:
                        try:
                            slide.shapes._spTree.remove(s._element)
                        except Exception:
                            pass
                        break

            # If only 1 photo: expand its slot to fill the combined width
            if num_valid == 1 and all_pic_shapes:
                path, (pic_name, txt_name) = valid_pairs[0]
                pic_s = all_pic_shapes.get(pic_name)
                txt_s = all_txt_shapes.get(txt_name)
                if pic_s and len(all_pic_shapes) > 1:
                    # Calculate total width spanned by all original slots
                    lefts = [s.left for s in all_pic_shapes.values()]
                    rights = [s.left + s.width for s in all_pic_shapes.values()]
                    total_left = min(lefts)
                    total_right = max(rights)
                    new_width = total_right - total_left
                    pic_s.left = total_left
                    pic_s.width = new_width
                    if txt_s:
                        txt_s.left = total_left
                        txt_s.width = new_width

            elif num_valid == 2 and all_pic_shapes:
                # Redistribute 2 photos evenly across the total span
                used_shapes = [(path, all_pic_shapes.get(cfg[0]), all_txt_shapes.get(cfg[1]))
                               for path, cfg in valid_pairs]
                all_lefts = [s.left for s in all_pic_shapes.values()]
                all_rights = [s.left + s.width for s in all_pic_shapes.values()]
                total_left = min(all_lefts)
                total_right = max(all_rights)
                gap = Emu(91440 * 8)  # ~8pt gap in EMU
                slot_w = (total_right - total_left - gap) // 2
                for i, (path, pic_s, txt_s) in enumerate(used_shapes):
                    if pic_s:
                        new_left = total_left + i * (slot_w + gap)
                        pic_s.left = new_left
                        pic_s.width = slot_w
                        if txt_s:
                            txt_s.left = new_left
                            txt_s.width = slot_w

        # Now fill the valid slots with actual images
        for path, (pic_name, txt_name) in valid_pairs:
            self._fill_image_slot(slide, pic_name, txt_name, path)

    def _fill_image_slots_grid(self, slide, slot_configs: list, photo_paths: list):
        """
        Arrange 4 to 6 photos in a clean grid (2x2 or 2x3) within the bounding box
        of the original placeholders, preserving aspect ratio (contain fit) and avoiding overlaps.
        """
        from pptx.util import Emu
        import os

        # 1. Collect all PIC_ and TXT_ shapes to get the bounding box
        pic_shapes = []
        txt_shapes = []
        for (pic_name, txt_name) in slot_configs:
            for s in slide.shapes:
                if s.name == pic_name:
                    pic_shapes.append(s)
                elif s.name == txt_name:
                    txt_shapes.append(s)
                    
        if not pic_shapes:
            return
            
        # 2. Compute bounding box
        min_left = min(s.left for s in pic_shapes)
        min_top = min(s.top for s in pic_shapes)
        max_right = max(s.left + s.width for s in pic_shapes)
        max_bottom = max(s.top + s.height for s in pic_shapes)
        
        box_w = max_right - min_left
        box_h = max_bottom - min_top
        
        # 3. Remove all original PIC_ and TXT_ shapes
        for s in pic_shapes + txt_shapes:
            try:
                slide.shapes._spTree.remove(s._element)
            except Exception:
                pass
                
        # 4. Filter only valid photos
        valid_paths = [p for p in photo_paths if p and os.path.exists(p) and os.path.getsize(p) > 0]
        num_photos = len(valid_paths)
        if num_photos == 0:
            return
            
        # Determine grid size (rows, cols)
        if num_photos <= 4:
            rows, cols = 2, 2
        else:
            rows, cols = 2, 3
            
        # Calculate cell dimensions (including a small gap, say ~6pt in EMU)
        gap_w = Emu(91440 * 6)  # ~6pt gap
        gap_h = Emu(91440 * 6)
        
        cell_w = (box_w - gap_w * (cols - 1)) // cols
        cell_h = (box_h - gap_h * (rows - 1)) // rows
        
        # 5. Insert images into the grid
        for i, path in enumerate(valid_paths):
            if i >= rows * cols:
                break
                
            r = i // cols
            c = i % cols
            
            cell_left = min_left + c * (cell_w + gap_w)
            cell_top = min_top + r * (cell_h + gap_h)
            
            # Process and fit image
            temp_dir = os.path.dirname(path)
            base = os.path.basename(path)
            fitted_path = os.path.join(temp_dir, f"fitted_grid_{cell_w}_{cell_h}_{base}")
            
            # Proportional sizing for image processor
            target_w = 600
            target_h = int(600 * cell_h / cell_w) if cell_w > 0 else 450
            
            try:
                ImageProcessor.process_and_fit_image(path, fitted_path, target_width=target_w, target_height=target_h)
                slide.shapes.add_picture(fitted_path, cell_left, cell_top, cell_w, cell_h)
            except Exception as e:
                print(f"Error drawing grid image: {e}")
                
            # Optionally add a small label at the bottom left of the cell
            try:
                from pptx.util import Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
                
                # Small transparent textbox
                tb = slide.shapes.add_textbox(cell_left, cell_top + cell_h - Emu(91440 * 12), cell_w, Emu(91440 * 12))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
                p = tf.paragraphs[0]
                p.text = f"Ảnh {i+1}"
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    run = p.runs[0]
                    run.font.name = FONT_PRIMARY
                    run.font.size = Pt(9)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception as e:
                print(f"Error adding text label in grid: {e}")

    def _replace_with_picture(self, slide, shape_name: str, img_path: str):
        """Find OLE or shape, delete, insert picture preserving z-order."""
        for shape in slide.shapes:
            if shape.name == shape_name:
                left, top, w, h = shape.left, shape.top, shape.width, shape.height
                
                spTree = slide.shapes._spTree
                placeholder_element = shape._element
                insert_index = spTree.index(placeholder_element)
                
                new_picture = slide.shapes.add_picture(img_path, left, top, w, h)
                new_pic_element = new_picture._element
                
                spTree.remove(new_pic_element)
                spTree.insert(insert_index, new_pic_element)
                spTree.remove(placeholder_element)
                return

    def _set_slide_title(self, slide, original_title_start: str, new_title: str):
        """Find title shape starting with original_title_start and replace its text, styling it as Arial bold 24pt."""
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip().startswith(original_title_start):
                tf = s.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = new_title
                p.font.name = "Arial"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                return

    def _highlight_rating(self, slide, rating_val: str):
        colors = {
            "tốt": RGBColor(226, 240, 217),      # light green
            "đạt": RGBColor(255, 242, 204),      # light yellow
            "chưa đạt": RGBColor(240, 206, 206)   # light red
        }
        selected_color = colors.get(rating_val.lower(), RGBColor(200, 200, 200))
        target_name = f"SHP_RATING_{rating_val.upper().replace(' ', '')}"
        
        for shape in slide.shapes:
            if shape.name == target_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = selected_color
                break

    def _replace_with_table(self, slide, anchor_name: str, rows_data: list, col_widths: list, headers: list, aligns: list):
        """Delete anchor shape and insert custom formatted native table."""
        anchor_shape = None
        for shape in slide.shapes:
            if shape.name == anchor_name:
                anchor_shape = shape
                break
                
        if anchor_shape is None:
            print(f"Warning: Table anchor '{anchor_name}' not found.")
            return

        left = anchor_shape.left
        top = anchor_shape.top
        width = anchor_shape.width
        height = anchor_shape.height
        
        slide.shapes._spTree.remove(anchor_shape._element)
        
        num_rows = len(rows_data) + 1
        num_cols = len(headers)
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table
        
        for col_idx, w in enumerate(col_widths):
            if col_idx < len(table.columns):
                table.columns[col_idx].width = w

        header_bg = RGBColor(10, 35, 66)
        
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            self._format_cell(cell, h, font_size=12, bold=True, color=RGBColor(255, 255, 255), align="center", bg_color=header_bg)
            
        for r_idx, row_vals in enumerate(rows_data):
            row_bg = RGBColor(255, 255, 255) if r_idx % 2 == 0 else RGBColor(245, 248, 252)
            for c_idx, val in enumerate(row_vals):
                if c_idx < num_cols:
                    cell = table.cell(r_idx + 1, c_idx)
                    align = aligns[c_idx]
                    bold = True if c_idx == 1 else False
                    self._format_cell(cell, val, font_size=11, bold=bold, color=RGBColor(0, 0, 0), align=align, bg_color=row_bg)

    def _format_cell(self, cell, text: str, font_size: int, bold: bool, color: RGBColor, align: str, bg_color: RGBColor):
        # Apply margins and word wrap
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        
        lines = str(text).split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = FONT_PRIMARY
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    def _set_cell_text_preserve_format(self, cell, text):
        """Write text to cell, preserving font name, size, bold, color, and paragraph alignment from cell's first run."""
        text = str(text)
        font_name = FONT_PRIMARY
        font_size = Pt(11)
        bold = False
        color = RGBColor(0, 0, 0)
        align = PP_ALIGN.LEFT
        
        tf = cell.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            align = p.alignment
            if p.runs:
                run = p.runs[0]
                if run.font.name:
                    font_name = run.font.name
                if run.font.size:
                    font_size = run.font.size
                bold = run.font.bold
                try:
                    if run.font.color and run.font.color.rgb:
                        color = run.font.color.rgb
                except AttributeError:
                    pass
                    
        # Apply margins and word wrap
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        tf.clear()
        
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = align
            
            p.font.name = font_name
            p.font.size = font_size
            p.font.bold = bold
            p.font.color.rgb = color
            
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _draw_product_grid(self, slide, anchor_name: str, items: list, is_best_seller: bool):
        """Draw a 2x5 grid of product images and text information instead of a table."""
        anchor_shape = None
        for shape in slide.shapes:
            if shape.name == anchor_name:
                anchor_shape = shape
                break
                
        if anchor_shape is None:
            print(f"Warning: Grid anchor '{anchor_name}' not found.")
            return

        # Position of anchor
        left = anchor_shape.left
        top = anchor_shape.top
        
        # Remove anchor shape
        slide.shapes._spTree.remove(anchor_shape._element)
        
        # Grid parameters (2 rows, 5 columns)
        rows = 2
        cols = 5
        
        # Let's define the total available width and height.
        # Since slide width is 20 inches, we can use 18.0 inches for content width.
        # Slide height is 11.25 inches. We can use 7.5 inches for content height.
        total_width = Inches(18.0)
        total_height = Inches(7.5)
        
        # Calculate cell widths and heights
        gap_x = Inches(0.2)
        gap_y = Inches(0.3)
        
        cell_width = (total_width - (cols - 1) * gap_x) / cols
        cell_height = (total_height - (rows - 1) * gap_y) / rows
        
        # Image placeholder dimensions (square)
        img_size = min(cell_width, cell_height - Inches(1.0))
        img_left_offset = (cell_width - img_size) / 2
        
        for idx in range(10):
            if idx >= len(items):
                break
            item = items[idx]
            
            # Row and column index
            r_idx = idx // cols
            c_idx = idx % cols
            
            # Calculate coordinates
            cell_left = left + c_idx * (cell_width + gap_x)
            cell_top = top + r_idx * (cell_height + gap_y)
            
            # 1. Draw Image Placeholder Shape
            img_left = cell_left + img_left_offset
            img_top = cell_top
            
            img_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, img_left, img_top, img_size, img_size
            )
            img_shape.fill.solid()
            img_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
            img_shape.line.color.rgb = RGBColor(200, 200, 200)
            img_shape.line.width = Pt(1)
            
            # Text inside image placeholder
            img_shape.text = "Ảnh sản phẩm\n(Tỷ lệ 1:1)"
            p_img = img_shape.text_frame.paragraphs[0]
            p_img.alignment = PP_ALIGN.CENTER
            p_img.font.name = FONT_PRIMARY
            p_img.font.size = Pt(9)
            p_img.font.color.rgb = RGBColor(0, 0, 0)
            img_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 2. Draw Text Box Below Image
            text_left = cell_left
            text_top = cell_top + img_size + Inches(0.08)
            text_width = cell_width
            text_height = cell_height - img_size - Inches(0.08)
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            # Row 1: STT and SKU
            p1 = tf.paragraphs[0]
            p1.text = f"STT {item.rank}. SKU: {item.sku}"
            p1.font.name = FONT_PRIMARY
            p1.font.size = Pt(9.5)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(10, 35, 66)
            p1.space_after = Pt(2)
            
            # Row 2: Product Name
            p2 = tf.add_paragraph()
            p2.text = item.product_name
            p2.font.name = FONT_PRIMARY
            p2.font.size = Pt(9)
            p2.font.color.rgb = RGBColor(0, 0, 0)
            p2.space_after = Pt(2)
            
            # Row 3: Details (Brand, Sales, Stock)
            p3 = tf.add_paragraph()
            if is_best_seller:
                p3.text = f"Hiệu: {item.brand}\nBán: {item.sales_4w:,.0f} | Tồn: {item.stock_qty:,.0f}"
            else:
                p3.text = f"Hiệu: {item.brand}\nTồn: {item.stock_qty:,.0f} | Tuổi: {item.age_days} ngày"
            p3.font.name = FONT_PRIMARY
            p3.font.size = Pt(8.5)
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(0, 0, 0)

    def _remove_shape_by_name(self, slide, name: str):
        """Helper to remove a shape by its exact name."""
        for s in list(slide.shapes):
            if s.name == name:
                try:
                    slide.shapes._spTree.remove(s._element)
                except Exception as e:
                    print(f"Error removing shape {name}: {e}")

    def _delete_slide(self, prs, slide):
        """Remove a slide from the presentation."""
        try:
            for s_id in prs.slides._sldIdLst:
                if s_id.id == slide.slide_id:
                    prs.slides._sldIdLst.remove(s_id)
                    break
        except Exception as e:
            print(f"Error deleting slide: {e}")

    def _cleanup_unused_placeholders(self, prs):
        """Delete any unfilled shapes starting with PIC_ or TXT_..._PLACEHOLDER / TXT_..._IMAGE_PLACEHOLDER."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for slide in prs.slides:
            # 1. Collect all modified text placeholders that should NOT be deleted
            modified_txt_names = set()
            for shape in slide.shapes:
                if shape.name.startswith("TXT_") and ("PLACEHOLDER" in shape.name or "IMAGE" in shape.name):
                    txt = shape.text_frame.text.strip() if shape.has_text_frame else ""
                    # Keep if modified (not empty, and doesn't contain camera prompt, placeholder or brackets)
                    if txt and "chèn ảnh" not in txt.lower() and "📷" not in txt and "nhập đề xuất" not in txt.lower() and "[nhập" not in txt.lower() and txt != "[Mô tả vấn đề]":
                        modified_txt_names.add(shape.name)
            
            shapes_to_delete = []
            for shape in slide.shapes:
                if shape.name.startswith("PIC_") and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    # Map PIC to corresponding TXT name
                    txt_name = shape.name.replace("PIC_", "TXT_")
                    if "_FRONTAGE_" in shape.name:
                        txt_name = shape.name.replace("PIC_FRONTAGE_", "TXT_FRONTAGE_IMAGE_PLACEHOLDER_")
                    elif "_INNER_" in shape.name:
                        txt_name = shape.name.replace("PIC_INNER_", "TXT_INNER_IMAGE_PLACEHOLDER_")
                    elif "_VM_" in shape.name:
                        txt_name = shape.name.replace("PIC_VM_", "TXT_VM_IMAGE_PLACEHOLDER_")
                    elif "_COMP_" in shape.name:
                        txt_name = shape.name.replace("PIC_COMP_", "TXT_COMP_IMAGE_PLACEHOLDER_")
                    elif shape.name == "PIC_STOCKROOM":
                        txt_name = "TXT_STOCKROOM_IMAGE_PLACEHOLDER"
                    elif shape.name == "PIC_FITTING_ROOM":
                        txt_name = "TXT_FITTING_ROOM_IMAGE_PLACEHOLDER"
                    elif shape.name == "PIC_CASHIER":
                        txt_name = "TXT_CASHIER_IMAGE_PLACEHOLDER"
                    elif shape.name == "PIC_CSVC_ISSUE":
                        txt_name = "TXT_CSVC_IMAGE_PLACEHOLDER"
                        
                    # If corresponding text was modified (e.g. set to 'Không có ảnh ghi nhận'), keep both the PIC frame and the text
                    if txt_name in modified_txt_names:
                        continue
                    shapes_to_delete.append(shape)
                    
                elif shape.name.startswith("TXT_") and ("PLACEHOLDER" in shape.name or "IMAGE" in shape.name):
                    # Keep if modified
                    if shape.name in modified_txt_names:
                        continue
                    shapes_to_delete.append(shape)
                    
            for shape in shapes_to_delete:
                try:
                    slide.shapes._spTree.remove(shape._element)
                except Exception as e:
                    print(f"Error removing unused shape {shape.name}: {e}")

    def _collect_resolved_comments(self, checklist_data, section_keys):
        """
        Collect resolved and pending items from the checklist data for specified section keys.
        Returns a dictionary with 'failures' (pending) and 'resolved' (resolved) lists of strings.
        """
        result = {"failures": [], "resolved": []}
        if not checklist_data or "sections" not in checklist_data:
            return result
            
        sections = checklist_data.get("sections", {})
        for sec_key in section_keys:
            sec_val = sections.get(sec_key, {})
            items = sec_val.get("items", [])
            for item in items:
                if item.get("eval") == "Không đạt":
                    item_label = item.get("label", "")
                    note = item.get("note", "Lỗi phát hiện")
                    resolved = item.get("resolved") or "Không"
                    severity = item.get("severity", "Trung bình")
                    assignee = item.get("assignee") or "CHT"
                    deadline = item.get("deadline") or "-"
                    
                    if resolved == "Có":
                        result["resolved"].append(f"- {item_label}: {note} (Đã khắc phục tại chỗ)")
                    else:
                        result["failures"].append(f"- {item_label}: {note} ({severity}) -> {assignee} xử lý, hạn {deadline}")
        return result
