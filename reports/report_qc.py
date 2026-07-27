import os
import zipfile
import re
import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

class QCViolationError(Exception):
    def __init__(self, message, violations=None):
        super().__init__(message)
        self.violations = violations or []

def normalize_qc_text(text: str) -> str:
    """Standardize ellipsis and whitespaces for robust matching."""
    if not text:
        return ""
    text = text.replace("…", "...")
    # Normalize whitespace
    text = " ".join(text.split())
    return text

def check_invalid_null(text: str):
    """Detect None, nan, NaN, null as independent values using word boundaries with context check."""
    # Find word matches for None, nan, NaN, null
    for m in re.finditer(r"\b(None|nan|NaN|null)\b", text):
        val = m.group(1)
        if val.lower() == "nan":
            start = max(0, m.start() - 10)
            end = min(len(text), m.end() + 10)
            context = text[start:end].lower()
            # If it's part of "gian nan" or "nan giải" or "nan giai", ignore it
            if "gian nan" in context or "nan giải" in context or "nan giai" in context:
                continue
        return True, val
    return False, None

class ReportQC:
    def __init__(self, config_path: str = "config/app_config.yaml"):
        self.root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.config_path = os.path.join(self.root_dir, config_path)
        with open(self.config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)
            
        # Whitelisted tokens allowed to remain
        self.whitelist_tokens = [
            r"\[/10\]",
            r"\[x\]",
            r"\[o\]"
        ]

    def verify_report(self, pptx_path: str, pdf_path: str, expected_slide_count: int, com_engine=None):
        """
        Runs the full 3-layer QC pipeline on the generated PowerPoint and PDF files.
        """
        print(f"Running quality control on generated files: {os.path.basename(pptx_path)}")
        
        # --- LAYER 1: ZIP INTEGRITY ---
        print("QC Layer 1: Checking ZIP integrity...")
        if not zipfile.is_zipfile(pptx_path):
            raise QCViolationError("File is not a valid zip archive (PowerPoint files must be zip format)")
            
        with zipfile.ZipFile(pptx_path) as z:
            corrupted = z.testzip()
            if corrupted is not None:
                raise QCViolationError(f"ZIP integrity check failed. Corrupted file inside ZIP: {corrupted}")
                
        # --- LAYER 2: BASIC XML PARSE & SLIDE COUNT ---
        print("QC Layer 2: Checking basic XML parse & slide count...")
        try:
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            if slide_count != expected_slide_count:
                raise QCViolationError(
                    f"Slide count discrepancy in PowerPoint!\n"
                    f"Expected: {expected_slide_count}, actual slide count: {slide_count}"
                )
        except Exception as e:
            if isinstance(e, QCViolationError):
                raise
            raise QCViolationError(f"XML parsing failed via python-pptx: {str(e)}")

        # --- LAYER 3: COM VERIFICATION ---
        if com_engine and com_engine.powerpoint:
            print("QC Layer 3: COM session active, slide count verified safely in Layer 2.")

        # --- PDF VERIFICATION ---
        print("QC PDF: Checking PDF page count...")
        if os.path.exists(pdf_path):
            try:
                reader = PdfReader(pdf_path)
                pdf_page_count = len(reader.pages)
                if pdf_page_count != expected_slide_count:
                    raise QCViolationError(
                        f"PDF page count ({pdf_page_count}) does not match PowerPoint slide count ({expected_slide_count})!"
                    )
            except Exception as e:
                if isinstance(e, QCViolationError):
                    raise
                raise QCViolationError(f"PDF structure is invalid or unreadable: {str(e)}")
        else:
            raise QCViolationError(f"PDF file was not exported or is missing at: {pdf_path}")

        # --- CONTENT QC: TOKEN SCAN ---
        print("QC Content: Scanning for unresolved placeholders and invalid values...")
        violations = []
        token_pattern = re.compile(r"\[[^\]]+\]")
        
        # Automation phrases to reject outside/inside brackets (case-insensitive normalized)
        blocked_phrases = [
            "nhận xét tổng quan đầu buổi kiểm tra...",
            "nhập đề xuất tại đây...",
            "nhập đề xuất phát triển tại đây...",
            "chèn ảnh tại đây",
            "chèn ảnh"
        ]

        for slide_idx, slide in enumerate(prs.slides):
            slide_id = ""
            for s in slide.shapes:
                if s.name == "META_SLIDE_ID" and s.has_text_frame:
                    slide_id = s.text_frame.text.strip()
                    break
            
            # Recursive search for all shapes including inside groups
            all_shapes = self._get_shapes_recursive(slide.shapes)
            
            for shape in all_shapes:
                # 1. Scan Text Frame
                if shape.has_text_frame:
                    self._check_text_field(
                        text=shape.text_frame.text,
                        slide_idx=slide_idx,
                        slide_id=slide_id,
                        shape_name=shape.name,
                        shape_id=shape.shape_id,
                        token_pattern=token_pattern,
                        blocked_phrases=blocked_phrases,
                        violations=violations
                    )
                
                # 2. Scan Table Cells
                if shape.has_table:
                    for r_idx, row in enumerate(shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            self._check_text_field(
                                text=cell.text_frame.text,
                                slide_idx=slide_idx,
                                slide_id=slide_id,
                                shape_name=shape.name,
                                shape_id=shape.shape_id,
                                token_pattern=token_pattern,
                                blocked_phrases=blocked_phrases,
                                violations=violations,
                                row_idx=r_idx,
                                col_idx=c_idx
                            )

        if violations:
            msg_parts = []
            for v in violations:
                loc = f"Slide index {v['slide_index']} (Slide {v['display_slide_number']}, ID: {v['slide_id']}), Shape '{v['shape_name']}' (ID: {v['shape_id']})"
                if v['row'] is not None:
                    loc += f", Cell [row {v['row']}, col {v['col']}]"
                msg_parts.append(f"- {loc}: Matched value '{v['matched_value']}' in context: '{v['text_context']}'")
            raise QCViolationError("QC Violations found:\n" + "\n".join(msg_parts), violations=violations)
            
        print("QC Pipeline completed successfully. Report is valid!")

    def _get_shapes_recursive(self, shape_collection):
        """Recursively collect shapes, expanding groups."""
        shapes = []
        for s in shape_collection:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    shapes.extend(self._get_shapes_recursive(s.shapes))
                except Exception as e:
                    print(f"Warning: could not recursively read group shapes: {e}")
            else:
                shapes.append(s)
        return shapes

    def _check_text_field(self, text, slide_idx, slide_id, shape_name, shape_id, token_pattern, blocked_phrases, violations, row_idx=None, col_idx=None):
        if not text:
            return
            
        normalized = normalize_qc_text(text)
        normalized_lower = normalized.lower()
        
        # 1. Reject camera icon "📷"
        if "📷" in text:
            violations.append({
                "slide_index": slide_idx,
                "display_slide_number": slide_idx + 1,
                "slide_id": slide_id,
                "shape_name": shape_name,
                "shape_id": shape_id,
                "row": row_idx,
                "col": col_idx,
                "matched_value": "📷",
                "text_context": text.strip()
            })
            return

        # 2. Check for blocked placeholder phrases
        for phrase in blocked_phrases:
            if phrase in normalized_lower:
                violations.append({
                    "slide_index": slide_idx,
                    "display_slide_number": slide_idx + 1,
                    "slide_id": slide_id,
                    "shape_name": shape_name,
                    "shape_id": shape_id,
                    "row": row_idx,
                    "col": col_idx,
                    "matched_value": phrase,
                    "text_context": text.strip()
                })
                return

        # 3. Check for invalid null/nan values
        has_null, null_val = check_invalid_null(text)
        if has_null:
            violations.append({
                "slide_index": slide_idx,
                "display_slide_number": slide_idx + 1,
                "slide_id": slide_id,
                "shape_name": shape_name,
                "shape_id": shape_id,
                "row": row_idx,
                "col": col_idx,
                "matched_value": null_val,
                "text_context": text.strip()
            })
            return

        # 4. Check for general unresolved bracket placeholders like [Tên cửa hàng]
        matches = token_pattern.findall(text)
        for m in matches:
            # Check if whitelisted
            is_whitelisted = False
            for wp in self.whitelist_tokens:
                if re.match(wp, m):
                    is_whitelisted = True
                    break
            if not is_whitelisted:
                violations.append({
                    "slide_index": slide_idx,
                    "display_slide_number": slide_idx + 1,
                    "slide_id": slide_id,
                    "shape_name": shape_name,
                    "shape_id": shape_id,
                    "row": row_idx,
                    "col": col_idx,
                    "matched_value": m,
                    "text_context": text.strip()
                })
                return
